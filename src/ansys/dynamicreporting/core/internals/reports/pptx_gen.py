#
# *************************************************************
#  Copyright 2022-2023 ANSYS, Inc.
#  All Rights Reserved.
#
#  Restricted Rights Legend
#
#  Use, duplication, or disclosure of this
#  software and its documentation by the
#  Government is subject to restrictions as
#  set forth in subdivision [(b)(3)(ii)] of
#  the Rights in Technical Data and Computer
#  Software clause at 52.227-7013.
# *************************************************************
#

import copy
import datetime
import itertools
import pickle
import uuid
from html.parser import HTMLParser as BaseHTMLParser
from io import BytesIO
from itertools import zip_longest
from pathlib import Path
from xml.sax.saxutils import escape, unescape

import numpy
from ceireports.utils import IntEnum, StrEnum
from data.extremely_ugly_hacks import safe_unpickle
from data.templatetags.data_tags import format_value_general
from data.utils import decode_table_data
from django.conf import settings
from django.template import engines
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.xmlchemy import OxmlElement
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.shapes.placeholder import TablePlaceholder
from pptx.util import Pt

from .engine import TemplateEngine


class PPTXManager:
    """
    Wrapper to manage a pptx file.
    """

    def __init__(self, input_filename):
        # load file
        try:
            self._presentation = Presentation(input_filename)
        except Exception as e:
            raise Exception(f"Unable to read file '{input_filename}': {e}")
        self._presentation_stream = None

    @property
    def presentation(self):
        return self._presentation

    @property
    def slides(self):
        return self._presentation.slides

    def get_slide(self, slide_id):
        return self._presentation.slides.get(int(slide_id))

    def get_presentation_stream(self):
        return self._presentation_stream.getvalue()

    @staticmethod
    def get_current_slide(placeholder):
        # SlidePlaceholder -> SlidePlaceholders -> Slide
        return placeholder._parent._parent

    @staticmethod
    def add_hyperlink_to_slide(run_object, source_slide, target_slide):
        """
        Hack to insert a hyperlink from a run object(paragraph) in a slide
        to another slide.
        """
        r_id = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
        r_pr = run_object._r.get_or_add_rPr()
        hyperlink_click = r_pr.add_hlinkClick(r_id)
        hyperlink_click.set('action', 'ppaction://hlinksldjump')

    @staticmethod
    def remove_placeholder(placeholder):
        # remove ph
        ph_elem = placeholder.element
        ph_elem.getparent().remove(ph_elem)

    @staticmethod
    def truncate_table(table, num_rows, num_cols):
        # CAVEAT: This hack to remove rows/cols can only be used for empty tables
        # or tables without links to other content.
        # https://github.com/scanny/python-pptx/pull/399#issuecomment-414149996
        while len(table.rows) > num_rows:
            row = table.rows[len(table.rows) - 1]
            table._tbl.remove(row._tr)
        while len(table.columns) > num_cols:
            column = table.columns[len(table.columns) - 1]
            col_idx = table._tbl.tblGrid.index(column._gridCol)
            for tr in table._tbl.tr_lst:
                tr.remove(tr.tc_lst[col_idx])
            table._tbl.tblGrid.remove(column._gridCol)

    @staticmethod
    def set_log_base(value_axis, base):
        # Inspired by:
        # https://github.com/scanny/python-pptx/issues/349#issuecomment-631261279
        # All it does is find a specific block to set the log base
        # and then modifies it to add the logBase block.
        scaling_elements = value_axis._element.xpath(r"c:scaling")
        if scaling_elements:
            scaling_element = scaling_elements[0]
        else:
            scaling_element = OxmlElement(r"c:scaling")
            value_axis._element.append(scaling_element)
        # now look for the target
        log_elements = scaling_element.xpath(r"c:logBase")
        if log_elements:
            log_element = log_elements[0]
            log_element.set("val", base)
        else:
            log_element = OxmlElement(r"c:logBase")
            log_element.set("val", base)
            value_axis._element.xpath(r"c:scaling")[0].append(log_element)

    def delete_slide(self, slide):
        """
        Delete a slide from presentation.

        python-pptx has no way to directly delete, so we need this awful hack.
        https://github.com/scanny/python-pptx/issues/67

        :param slide: slide object to delete
        :return:
        """
        slide_rel_map = {}
        for i, slide_obj in enumerate(self._presentation.slides._sldIdLst):
            slide_rel_map[slide_obj.id] = (i, slide_obj.rId)
        # fetch
        idx, rel_idx = slide_rel_map[slide.slide_id]
        # drop relation to slide's part
        self._presentation.part.drop_rel(rel_idx)
        # delete reference from slide list
        del self._presentation.slides._sldIdLst[idx]

    def duplicate_slide(self, source_slide, target_idx):
        """
        Another awful hack to duplicate a slide and insert at a specified index.

        "Inspired" by:
         https://github.com/scanny/python-pptx/issues/132
         https://github.com/scanny/python-pptx/issues/754
         https://github.com/scanny/python-pptx/issues/68
         https://github.com/scanny/python-pptx/issues/274
         https://stackoverflow.com/a/51154583/3628052

        NOTE/WARNING: Only works for versions of python-pptx <= 0.6.19
        due to an issue with slide.part.rels and slide.part.rels.add_relationship

        :param source_slide:  the source slide
        :param target_idx: the target position in the slide element list
        :return: copied slide
        """
        # append a slide with source's layout
        copied_slide = self._presentation.slides.add_slide(source_slide.slide_layout)
        # delete shapes to get a blank slide
        for shape in copied_slide.shapes:
            shape.element.getparent().remove(shape.element)
        # copy shapes from source
        for shape in source_slide.shapes:
            copied_shape = copy.deepcopy(shape.element)
            copied_slide.shapes._spTree.insert_element_before(copied_shape, 'p:extLst')
        # copy rels from source
        for key, value in source_slide.part.rels.items():
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                target = value._target
                # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
                if "chart" in value.reltype:
                    partname = target.package.next_partname(
                        ChartPart.partname_template)
                    xlsx_blob = target.chart_workbook.xlsx_part.blob
                    target = ChartPart(partname, target.content_type,
                                       copy.deepcopy(target._element), package=target.package)
                    target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(
                        xlsx_blob, target.package)
                # copy relationship
                copied_slide.part.rels.add_relationship(value.reltype, target, value.rId, value.is_external)
        # insert copied slide at the target index
        # lxml does not make copies unless you tell it to.
        # If you stick an element somewhere else, it comes out of wherever it was before.
        self._presentation.slides.element.insert(target_idx, self._presentation.slides.element[-1])
        return copied_slide

    def save_presentation(self):
        # https://stackoverflow.com/a/42142698/3628052
        # TLDR; If you don't re-instantiate every time, you will
        # get a corrupt file.
        self._presentation_stream = BytesIO()
        self._presentation.save(self._presentation_stream)
        self._presentation = Presentation(self._presentation_stream)


class NoContentError(Exception):
    pass


class UnsupportedTagsError(Exception):
    pass


class NestedTagsError(Exception):
    pass


class HTMLParser(BaseHTMLParser):
    """
    Very naive parser that satisfies the PPTX use case only.
    Supports a combination of non-nested sequence of HTML tags,
    with or without content and possibly some plain text.
    """
    ALLOWED_VOID_TAGS = ["br"]  # https://html.spec.whatwg.org/#void-elements
    ALLOWED_HEADER_TAGS = ["h1", "h2", "h3", "h4", "h5", "h6"]
    ALLOWED_OTHER_TAGS = ["a"]
    ALLOWED_TAGS = ALLOWED_VOID_TAGS + ALLOWED_HEADER_TAGS + ALLOWED_OTHER_TAGS

    def __init__(self):
        super().__init__()
        self._start_tags = []
        self._data = []
        self._current_tag = None

    @property
    def tags(self):
        return self._start_tags

    def handle_starttag(self, tag, attrs):
        if tag not in self.ALLOWED_TAGS:
            raise UnsupportedTagsError(f"The HTML tag '{tag}' is not supported.")
        if self._current_tag is None:
            self._current_tag = tag
        else:
            raise NestedTagsError(f"The tag '{tag}' is nested within the tag '{self._current_tag}'.")
        self._start_tags.append(tag)
        # void elements do not have any data contained in them, so we use fillers.
        if tag in self.ALLOWED_VOID_TAGS:
            self._data.append("")
            # Void tags can be represented as <tag> or <tag />.
            # In the former, handle_endtag is never called, so it
            # is never closed. Just close it early here.
            self._current_tag = None

    def handle_endtag(self, tag):
        if tag == self._current_tag:
            self._current_tag = None

    def handle_data(self, data):
        # if data is encountered without any tags,
        # use filler for tag name.
        if self._current_tag is None:
            self._start_tags.append("")
        self._data.append(data)

    def parse_content(self, html, plaintext=False):
        self.feed(html)
        # return plaintext
        if plaintext:
            text = ""
            for data in self._data:
                text += f"{data} "
            return text.strip()
        # return structured data
        return list(zip(self._start_tags, self._data))

    def is_html(self, string):
        self.feed(string)
        for tag in self._start_tags:
            if tag:
                return True
        return False

    def reset(self):
        super().reset()
        self._start_tags = []
        self._data = []
        self._current_tag = None


class PPTXGenerator:
    """
    PPTX Report Generation
    """
    # PP_PLACEHOLDER are placeholder enums.
    # MSO_SHAPE_TYPE are enums to recognize a shape post-insertion
    # for example: an existing table.
    ITEM_PLACEHOLDER_MAP = {
        "string": [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY],
        "html": [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY],
        "image": [PP_PLACEHOLDER.PICTURE],
        "anim": [PP_PLACEHOLDER.MEDIA_CLIP],
        "table": [PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.CHART],
        "tree": [PP_PLACEHOLDER.TABLE],
        "scene": [PP_PLACEHOLDER.PICTURE],
        "file": [PP_PLACEHOLDER.PICTURE],
    }
    ALLOWED_PLACEHOLDERS = tuple(itertools.chain.from_iterable(ITEM_PLACEHOLDER_MAP.values()))
    ITEM_SPLIT_ALLOWED_TYPES = ("string", "table", "tree", "html")
    # obtained from https://getbootstrap.com/docs/5.0/content/typography/
    HTML_TAG_FONT_MAP = {
        "h1": 40,
        "h2": 32,
        "h3": 28,
        "h4": 24,
        "h5": 20,
        "h6": 16,
    }
    TOC_LINK_DEFAULT_TEXT = "Table of contents"
    PARAM_PREFIX = "Nexus:"

    # params
    class Param(StrEnum):
        TOC = "toc"
        TOC_LINK = "toc_link"
        HTML_HEADER = "html_header"
        COMMENTS = "comments"
        BREADCRUMB = "breadcrumbs"

    # defaults
    class DefaultLimit(IntEnum):
        LINE = 7
        ROW = 10
        COL = 5

    # link types
    class HyperlinkType(StrEnum):
        SLIDE = "slide"

    def __init__(self, pptx_mgr):
        # ppt
        self.pptx_mgr = pptx_mgr
        # slide id mapped to its info
        self.ppt_info = {}
        # slide id to number of dupes needed
        self.duplicate_count_info = {}
        # toc slide's id
        self.toc_slide_id = None
        # list of slide ids to exclude from toc
        self.toc_exclude_set = set()
        # slide_ids to titles
        self.title_info = {}
        # hierarchy
        self.hierarchy_info = {}

    def _add_string(self, placeholder, data):
        lines = data["data"]
        item_ctx = data["item_ctx"]
        num_lines = item_ctx["num_lines"]
        extra_info = item_ctx.get("extra_info")  # optional extra info for each line
        # treat each line as a run instead of a para
        as_run = item_ctx.get("as_run", False)
        # one para = one line of text
        # one paragraph already exists by default
        if not as_run:
            for _ in range(num_lines - 1):
                placeholder.text_frame.add_paragraph()
        # - A text placeholder is made of paragraphs. There is one empty paragraph
        # by default, all the time.
        # - A paragraph is made of one or more 'Runs', which are just streams of text.
        # - If a paragraph is empty(default), the para.runs tuple will be empty, so no Runs.
        # - If there is text, there will be a Run object.
        # - Both paragraphs and runs have a text attr that returns the contained text.
        # - If a para has only one run, then para.text == para.runs[0].text
        current_slide = self.pptx_mgr.get_current_slide(placeholder)
        for idx, line in enumerate(lines):
            run = None
            if as_run:
                para = placeholder.text_frame.paragraphs[0]
                # check for existing text and use it
                if len(para.runs) > idx:
                    run = para.runs[idx]
            else:
                para = placeholder.text_frame.paragraphs[idx]
                # check for existing text and use it
                if para.runs:
                    run = para.runs[0]
            if run is None:
                run = para.add_run()
                run.text = line
            # extras
            if extra_info is not None:
                line_num = idx
                # check if we are continuing from a previous part
                # and start at the appropriate index.
                part_num = item_ctx.get("part_num")
                if part_num is not None and part_num > 0:
                    part_limit = item_ctx["part_limit"]
                    line_num = idx + part_limit * part_num
                extra = extra_info.get(line_num)
                if extra is None:
                    continue
                # hyperlink for the line
                link_type = extra.get("link_type")
                link_target = extra.get("link_target")
                try:
                    if link_type == self.HyperlinkType.SLIDE:
                        target_slide = self.pptx_mgr.get_slide(link_target)
                        self.pptx_mgr.add_hyperlink_to_slide(run, current_slide, target_slide)
                except Exception as e:
                    raise Exception(f"There was an error creating hyperlinks for the text '{line}':: {e}")
                # depth for the line
                depth = extra.get("depth")
                if depth is not None:
                    para.level = depth

    def _add_html(self, placeholder, data):
        parsed_html = data["data"]
        add_new_para = False
        # take the default para as current
        current_para = placeholder.text_frame.paragraphs[0]
        for idx, tag_content in enumerate(parsed_html):
            tag, content = tag_content
            if idx == 0 or not add_new_para:
                para = current_para
            else:
                para = placeholder.text_frame.add_paragraph()
                current_para = para
            run = para.add_run()
            run.text = content
            if content:
                # set font size based on tag
                run.font.size = Pt(self.HTML_TAG_FONT_MAP.get(tag, 16))
            # when to break to a new line?
            if tag in HTMLParser.ALLOWED_HEADER_TAGS or tag == "br":
                add_new_para = True

    @staticmethod
    def _adjust_aspect_ratio(src_img, placeholder):
        """
        Find the original image's aspect ratio and compare it with the
        placeholder's(inserted content). If the ratio difference is greater
        than zero, it means the width of the placeholder is higher so we
        crop at the left and right ends. If the ratio difference is lesser
        than zero, it means the height of the placeholder is higher so we
        crop at the top and bottom ends. If both the height and width are
        higher, we only adjust the height to maintain the ratio. If both the
        height and width are lower, we only adjust the width to maintain the ratio.
        """
        if placeholder.width is None or placeholder.height is None:
            return

        ph_ratio = placeholder.width / placeholder.height
        img_width, img_height = src_img.size
        img_ratio = img_width / img_height
        ratio_difference = ph_ratio - img_ratio

        # Placeholder is wider
        if ratio_difference > 0:
            difference_at_ends = ratio_difference / 2
            placeholder.crop_left = -difference_at_ends
            placeholder.crop_right = -difference_at_ends
        else:
            # Placeholder is taller
            difference_at_ends = -ratio_difference / 2
            placeholder.crop_bottom = -difference_at_ends
            placeholder.crop_top = -difference_at_ends

    def _insert_image(self, placeholder, img_path):
        picture = placeholder.insert_picture(img_path)
        # revert auto-cropping initially
        crop_attrs = ["crop_bottom", "crop_top", "crop_right", "crop_left"]
        for attr in crop_attrs:
            setattr(picture, attr, 0.0)
        # adjust
        self._adjust_aspect_ratio(picture.image, picture)

    def _add_image(self, placeholder, data):
        img_path = data["payload_server_pathname"]
        if Path(img_path).is_file():
            self._insert_image(placeholder, img_path)

    def _add_anim(self, placeholder, data):
        # get details before removing
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
        name = placeholder.name
        # remove placeholder
        self.pptx_mgr.remove_placeholder(placeholder)
        slide = self.pptx_mgr.get_current_slide(placeholder)
        movie_filename = data["payload_server_pathname"]
        # generate thumbnail
        try:
            poster_image = Path(settings.MEDIA_ROOT).resolve(
                strict=True) / f'thumbnail_{data["guid"]}_{data["type"]}.png'
            # create only if required.
            if not Path.is_file(poster_image):
                import enve
                movie = enve.movie(enve.MOVIE_READ)
                movie.filename = movie_filename
                if movie.open() > -1:
                    image = movie.getframe(0)[0]
                    movie.close()
                    image.save(str(poster_image))
        except Exception:
            poster_image = None

        if not poster_image:
            raise Exception(f"Unable to find or generate the poster frame image needed to"
                            f" insert item with guid '{data['guid']}'")

        # add movie
        # We must use the size of the placeholder from before
        # and then adjust the content later after insertion.
        # There is no auto-cropping here like in picture insertion
        # so, we don't worry about that.
        movie = slide.shapes.add_movie(movie_filename,
                                       left, top, width, height,
                                       poster_frame_image=str(poster_image),
                                       mime_type='video/mp4')
        movie.name = name
        # adjust
        self._adjust_aspect_ratio(movie.poster_frame, movie)

    def _fill_table(self, placeholder, data):
        table_data = data["data"]
        item_ctx = data["item_ctx"]
        num_rows = item_ctx["num_rows"]
        num_cols = item_ctx["num_cols"]
        # get table
        if isinstance(placeholder, TablePlaceholder):
            graphic_frame = placeholder.insert_table(rows=num_rows, cols=num_cols)
            table = graphic_frame.table
            created = True
        else:  # isinstance(placeholder, PlaceholderGraphicFrame)
            table = placeholder.table
            # resize table
            if len(table.rows) > num_rows or len(table.columns) > num_cols:
                self.pptx_mgr.truncate_table(table, num_rows, num_cols)
            created = False
        # build
        for row_idx, row in enumerate(table_data):
            for col_idx, col in enumerate(row):
                table.cell(row_idx, col_idx).text = col
        return table, created

    def _add_table(self, placeholder, data):
        item_ctx = data["item_ctx"]
        table, created = self._fill_table(placeholder, data)
        # set props only if a table was inserted by us
        if created:
            table.first_row = item_ctx["header_first_row"]
            table.first_col = item_ctx["header_first_col"]
        # remove header formatting for subsequent parts.
        part_num = item_ctx.get("part_num")
        # continuing parts should not have column label formatting
        if part_num is not None and part_num > 0:
            table.first_row = False

    def _add_tree(self, placeholder, data):
        table, created = self._fill_table(placeholder, data)
        item_ctx = data["item_ctx"]
        part_num = item_ctx.get("part_num")
        # remove header formatting for subsequent parts or if we created it
        # if we created it OR the user did and it is a continuing part.
        if (part_num is not None and part_num > 0) or created:
            table.first_row = False
            table.first_col = False

    def _add_plot(self, placeholder, data):
        chart_data = data["data"]
        item_ctx = data["item_ctx"]
        chart_type = item_ctx["chart_type"]

        # insert chart
        graphic_frame = placeholder.insert_chart(chart_type, chart_data)
        chart = graphic_frame.chart

        # check for X/category axis
        try:
            x_axis = chart.category_axis
        except ValueError:
            pass
        else:
            if "x_title" in item_ctx:
                x_axis.axis_title.text_frame.text = item_ctx["x_title"]

        # check for Y/value axis
        try:
            y_axis = chart.value_axis
        except ValueError:
            pass
        else:
            if "y_title" in item_ctx:
                y_axis.axis_title.text_frame.text = item_ctx["y_title"]
            # log base
            if "log_base" in item_ctx:
                log_base = item_ctx["log_base"]
                try:
                    int(log_base)
                except ValueError as e:
                    raise Exception(f"The log base value '{log_base}' for the {chart_type} item "
                                    f"with guid \'{data['guid']}\' is not a valid integer.")
                # set
                self.pptx_mgr.set_log_base(y_axis, log_base)

        # plot title
        plot_title = item_ctx.get("plot_title")
        if plot_title is not None:
            chart.chart_title.text_frame.text = plot_title
        else:
            chart.has_title = False

        # legend
        chart.has_legend = bool(item_ctx["show_legend"])
        if chart.has_legend:
            chart.legend.font.size = Pt(10)
            # position
            chart.legend.position = getattr(XL_LEGEND_POSITION, item_ctx["legend_pos"].upper())
            # overlap legend with chart
            chart.legend.include_in_layout = bool(item_ctx["legend_overlap"])

        # data labels
        if chart_type == XL_CHART_TYPE.PIE:
            chart.plots[0].has_data_labels = True
            data_labels = chart.plots[0].data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.number_format = "0.00%"
            data_labels.font.size = Pt(8)
            data_labels.position = XL_LABEL_POSITION.INSIDE_END

    def _add_scene(self, placeholder, data):
        # check for proxy
        item_path = Path(data["payload_server_pathname"])
        proxy_path = item_path.parent / item_path.stem / "proxy.png"
        if Path(proxy_path).is_file():
            self._insert_image(placeholder, str(proxy_path))

    def _add_file(self, placeholder, data):
        item_path = Path(data["payload_server_pathname"])
        # check for proxy
        proxy_path = item_path.parent / item_path.stem / "proxy.png"
        if Path(proxy_path).is_file():
            self._insert_image(placeholder, str(proxy_path))

    def _add_item(self, placeholder, data):
        render_map = {
            "string": self._add_string,
            "html": self._add_html,
            "image": self._add_image,
            "anim": self._add_anim,
            "table": self._add_table,
            "tree": self._add_tree,
            "plot": self._add_plot,
            "scene": self._add_scene,
            "file": self._add_file
        }
        # get the 'processed' type
        item_type = data["type"]
        add_method = render_map[item_type]
        # call
        add_method(placeholder, data)

    def _verify_type_and_raise(self, placeholder, item):
        item_type = item.type
        if item_type not in self.ITEM_PLACEHOLDER_MAP:
            raise Exception(f"Support for the '{item_type}' type specified as part of the filter query"
                            f" in placeholder '{placeholder.name}' is currently not available.")
        allowed_placeholder_types = self.ITEM_PLACEHOLDER_MAP[item_type]
        ph_type = placeholder.placeholder_format.type
        if ph_type not in allowed_placeholder_types:
            raise Exception(f"The inserted placeholder '{placeholder.name}' does not match"
                            f" the type '{item_type}' of the item '{item.name}' returned by the filter query.")

    @staticmethod
    def _parse_params(input_string):
        parsed_dict = {}
        params_and_values = input_string.split("&")
        for p in params_and_values:
            param, _, value = p.strip().partition("=")
            parsed_dict[param.strip()] = value.strip()
        return parsed_dict

    @staticmethod
    def _process_string(item, ctx):
        payload = safe_unpickle(item.payloaddata)
        if not payload.strip():
            raise NoContentError
        lines = payload.splitlines()
        item_ctx = {
            "num_lines": len(lines),
            "extra_info": ctx.get("extra_info"),
            "as_run": ctx.get("as_run"),
        }
        title = ctx.get("title", False)
        if title and lines:
            # first line of valid text for toc title
            for line in lines:
                if line.strip():
                    item_ctx["title"] = line
                    break

        return {
            "type": "string",
            "data": lines,
            "item_ctx": item_ctx
        }

    @staticmethod
    def _process_html(item, ctx):
        from data.templatetags.data_tags import convert_macro_slashes
        payload = safe_unpickle(item.payloaddata)
        if not payload.strip():
            raise NoContentError
        try:
            converted_html = convert_macro_slashes(payload)
            django_engine = engines['django']
            django_template = django_engine.from_string('{% load data_tags %}' + converted_html)
            html = django_template.render(context=ctx)
            # parse
            parsed_html = HTMLParser().parse_content(html)
        except UnsupportedTagsError as e:
            raise Exception(f"The HTML item '{item.name}' with guid '{item.guid}' has HTML tags that are not supported."
                            f" Please check the PPTX export documentation for supported tags: {e}")
        except NestedTagsError as e:
            raise Exception(f"The HTML item '{item.name}' with guid '{item.guid}' has HTML tags that are nested."
                            f" Nesting of HTML tags is not supported in PPTX export at the moment: {e}")
        except Exception as e:
            raise Exception(f"Error processing HTML item '{item.name}' with guid '{item.guid}': {e}")

        item_ctx = {"num_lines": len(parsed_html)}
        title = ctx.get("title", False)
        if title and parsed_html:
            # first line of valid text for toc title
            for tag, text in parsed_html:
                if text.strip():
                    item_ctx["title"] = text
                    break

        return {
            "type": "html",
            "data": parsed_html,
            "item_ctx": item_ctx
        }

    @staticmethod
    def _process_table(item, ctx):
        table_data = safe_unpickle(item.payloaddata)

        # if the numpy array is empty, what's the point?
        if table_data["array"].size == 0:
            raise NoContentError

        # We have to decode each element in the numpy array explicitly to utf-8 strings before rendering
        # BUT only if the dtype is bytes.
        # Be aware that row and col labels can be bytes as well, as these are derived from table data.
        # so they have to be decoded as well.
        if table_data["array"].dtype.type == numpy.bytes_:
            table_data = decode_table_data(table_data)

        np_data = table_data["array"]
        plot_type = table_data.get('plot', 'table')
        num_rows, num_cols = np_data.shape

        # tables
        if plot_type == "table":
            col_labels = item.get_column_labels(table_data, ctx)
            row_labels = item.get_row_labels(table_data, ctx)

            # we add fillers so that we have column labels always and
            # the <thead> is rendered for all cases.
            if col_labels is None:
                # sometimes, if there's no column labels, the <thead> block is completely
                # ignored out of the generated DOM, outputting an incomplete HTML table.
                # so we add a 'filler' th to give a proper table.
                # This is now controlled by a property: table_default_col_labels
                table_default_col_labels = item.get_default(table_data, ctx,
                                                            'table_default_col_labels',
                                                            default=1,
                                                            force_int=True)
                if table_default_col_labels == 1:
                    col_labels = [f"Column-{i + 1}" for i in range(int(num_cols))]

            header_first_row = False
            header_first_col = False
            if row_labels:
                num_cols += 1
                header_first_col = True
            if col_labels:
                num_rows += 1
                header_first_row = True

            added_column_to_mask_rowlbls = False
            if row_labels is None:
                row_labels = [None] * num_rows
            else:
                # only if col lbls is an existing valid list
                if isinstance(col_labels, list):
                    # there are row labels giving us an extra column at the start
                    # so prepend '' to the column labels for a dummy column label.
                    col_labels.insert(0, '')
                    added_column_to_mask_rowlbls = True

            labels_and_rows = []
            # take filler column labels into account, but not filler row labels (ignore None's).
            if col_labels:
                for idx, label in enumerate(col_labels):
                    if label is not None:
                        if idx == 0 and added_column_to_mask_rowlbls:
                            cleaned_label = ""
                        else:
                            col_label_format = item.get_indexed_default(table_data,
                                                                        ctx,
                                                                        idx,
                                                                        'format_column',
                                                                        default='str')
                            cleaned_label, _ = item.format_table_value(label.strip(),
                                                                       col_label_format,
                                                                       context=ctx,
                                                                       item=item)
                    else:
                        cleaned_label = ""

                    if idx == 0:
                        labels_and_rows.append((cleaned_label, []))
                    else:
                        labels_and_rows[0][1].append(cleaned_label)

            # continue building final table data
            labels_and_rows.extend(list(zip(row_labels, np_data)))

            nan_display = item.get_default(table_data, ctx, 'nan_display', default='NaN')
            is_string = True
            is_float = False

            cleaned_data = []
            # initialize
            for _ in range(num_rows):
                cleaned_data.append([None] * num_cols)

            for row_idx, label_row in enumerate(labels_and_rows):
                label, row = label_row
                # set label if available
                if label is not None:
                    row_label_format = item.get_indexed_default(table_data, ctx, row_idx, 'format_row',
                                                                default='str')
                    formatted_label, _ = item.format_table_value(label.strip(),
                                                                 row_label_format,
                                                                 context=ctx,
                                                                 item=item)
                    cleaned_data[row_idx][0] = formatted_label
                # get dtype only if available. Some rows may not be numpy-formatted.
                if hasattr(row, "dtype"):
                    is_string = row.dtype.kind in "U"
                    is_float = row.dtype.kind in "f"
                # walk through columns and write to cells
                for col_idx, val in enumerate(row):
                    # if there was a label at (row_idx, 0), we start from the next position
                    real_col_idx = col_idx
                    if label is not None:
                        real_col_idx += 1
                    # use col_idx of the source data to retrieve properties and
                    # real_col_idx of the target table for insertion
                    if is_string:
                        number_format = item.get_indexed_default(table_data, ctx, col_idx, 'format', default='str')
                        needs_expansion = '{{' in val and '}}' in val
                        str_rep, _ = item.format_table_value(val.strip(), number_format, context=ctx,
                                                             item=item)
                        if needs_expansion:
                            str_rep = HTMLParser().parse_content(str_rep, plaintext=True)
                    elif is_float:
                        number_format = item.get_indexed_default(table_data, ctx, col_idx, 'format',
                                                                 default='scientific')
                        if numpy.isnan(val):
                            str_rep = nan_display
                        else:
                            str_rep, _ = format_value_general(val, number_format)
                    else:
                        str_rep = "Unknown:{}".format(np_data.dtype.kind)
                    cleaned_data[row_idx][real_col_idx] = str_rep

            payload = {
                "type": "table",
                "data": cleaned_data,
                "item_ctx": {
                    "num_rows": num_rows,
                    "num_cols": num_cols,
                    "header_first_row": header_first_row,
                    "header_first_col": header_first_col
                }
            }
        else:
            # Plots
            row_labels = item.get_row_labels(table_data, ctx, [None] * num_rows)
            col_labels = item.get_column_labels(table_data, ctx, list(range(num_cols)))

            from data.models import XAxisObj
            xaxis_obj = XAxisObj(ctx, table_data, row_labels, col_labels, item)

            y_data = []
            x_data = []
            y_labels = []
            for idx in xaxis_obj.y_row_indices():
                x_data.append(xaxis_obj.data(idx))
                y_data.append(np_data[idx])
                y_labels.append(row_labels[idx])

            chart_data = CategoryChartData()

            # only 3 supported plot types: line,bar,pie.
            if plot_type == "line":
                chart_type = XL_CHART_TYPE.LINE_MARKERS
            elif plot_type == "bar":
                chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
            elif plot_type == "pie":
                chart_type = XL_CHART_TYPE.PIE
            else:
                raise Exception(f"Plot type specified for {item.name} is not supported.")

            for idx in range(len(y_data)):
                # each iteration is one series
                input_data = []
                # if it is a pie, x values are constant
                if chart_type == XL_CHART_TYPE.PIE:
                    input_data.append(col_labels)
                else:
                    input_data.append(x_data[idx])
                input_data.append(y_data[idx])
                x_values, y_values = item.clean_plot_data(input_data)

                if chart_type in [XL_CHART_TYPE.LINE_MARKERS, XL_CHART_TYPE.COLUMN_CLUSTERED]:
                    if y_labels[idx] is not None:
                        series_label = item.clean_label(y_labels[idx])
                    else:
                        series_label = f"Row {idx}"
                else:  # pie
                    if row_labels[idx] is not None:
                        series_label = item.clean_label(row_labels[idx])
                    else:
                        series_label = f"Row {idx}"

                chart_data.categories = x_values
                chart_data.add_series(series_label, y_values)

            item_ctx = {"chart_type": chart_type}

            # titles
            # x axis
            x_title = xaxis_obj.title()  # reads the 'xtitle' property implicitly
            if x_title:
                item_ctx["x_title"] = x_title

            # y axis
            default_y_title = y_labels[0] if len(y_labels) == 1 else None
            y_title = ctx.get('ytitle', table_data.get('ytitle', default_y_title))
            if y_title:
                item_ctx["y_title"] = y_title

            # plot title
            plot_title = item.get_default(table_data, ctx, "plot_title", item.get_default(table_data, ctx, 'title'))
            if plot_title is not None:
                item_ctx["plot_title"] = item.clean_label(plot_title)

            # legend
            item_ctx["show_legend"] = item.get_default(table_data, ctx, "show_legend", default=1, force_int=True)
            # legend position
            item_ctx["legend_pos"] = item.get_default(table_data, ctx, "pptx_legend_position", default="corner")
            # legend overlap
            item_ctx["legend_overlap"] = item.get_default(table_data, ctx, "pptx_legend_overlap", default=0,
                                                          force_int=True)
            # log base
            log_base = item.get_default(table_data, ctx, "pptx_log_base", default=None)
            if log_base is not None:
                item_ctx["log_base"] = log_base

            payload = {
                "type": "plot",
                "data": chart_data,
                "item_ctx": item_ctx
            }

        return payload

    @staticmethod
    def _process_tree(item, ctx):
        def _get_formatted_value(input_val):
            formatted_val = str(input_val)  # handles a lot of cases...
            # formatting for the types: [float, int, datetime.datetime, str, bool, uuid.UUID, None]
            if input_val is None:
                formatted_val = ""
            elif type(input_val) == float:
                fmt = str(ctx.get('tree_format_float', 'floatdot2'))
                formatted_val = format_value_general(input_val, fmt)[0]
            elif type(input_val) == bool:
                fmt = str(ctx.get('tree_format_bool', 'True#False'))
                fmt = fmt.split('#')
                while len(fmt) < 2:
                    fmt.append("?")
                if input_val:
                    formatted_val = str(fmt[0])
                else:
                    formatted_val = str(fmt[1])
            elif type(input_val) == datetime.datetime:
                fmt = str(ctx.get('tree_format_date', 'date_44'))
                formatted_val = format_value_general(input_val, fmt)[0]
            if '{{' in formatted_val and '}}' in formatted_val:
                formatted_val = TemplateEngine.context_expansion(formatted_val, ctx, item=item)
                formatted_val = HTMLParser().parse_content(formatted_val, plaintext=True)
            return formatted_val

        def _build_tree(children, indent=""):
            rows = []
            for node in children:
                name = node.get('name', '')
                value = node.get('value', None)
                children = node.get('children', None)
                leaves = []
                if children:
                    name = f"- {name}"
                    leaves = _build_tree(children, indent + "      ")
                # support for multi-valued tree nodes.
                if isinstance(value, list):
                    r = [f"{indent}{name}", *[_get_formatted_value(v) for v in value]]
                else:
                    r = [f"{indent}{name}", _get_formatted_value(value)]
                rows.append(r)
                if leaves:
                    rows.extend(leaves)
            return rows

        def _get_dims(row_data):
            # avoid numpy here because rows are not always homogeneous
            n_cols = 0
            for r in row_data:
                if len(r) > n_cols:
                    n_cols = len(r)
            return len(row_data), n_cols

        data = safe_unpickle(item.payloaddata)
        if not data:
            raise NoContentError
        tree_data = _build_tree(data)
        num_rows, num_cols = _get_dims(tree_data)

        cleaned_data = []
        # initialize
        # use empty strings instead of None because
        # trees are not homogeneous and empty values
        # need to be displayed
        for _ in range(num_rows):
            cleaned_data.append([""] * num_cols)

        for row_idx, row in enumerate(tree_data):
            for col_idx, col in enumerate(row):
                cleaned_data[row_idx][col_idx] = str(col)

        return {
            "type": "tree",
            "data": cleaned_data,
            "item_ctx": {
                "num_rows": num_rows,
                "num_cols": num_cols
            }
        }

    def _process_item(self, item, ctx):
        # only certain types require processing
        try:
            if item.type == "string":
                processed_data = self._process_string(item, ctx)
            elif item.type == "html":
                processed_data = self._process_html(item, ctx)
            elif item.type == "table":
                processed_data = self._process_table(item, ctx)
            elif item.type == "tree":
                processed_data = self._process_tree(item, ctx)
            else:  # simple processing for the others
                processed_data = {
                    "type": item.type,
                    "data": None,
                    "item_ctx": {},
                    "payload_server_pathname": item.get_payload_server_pathname(),
                }
        except NoContentError:
            raise Exception(f"The {item.type} item '{item.name}' with guid '{item.guid} does not have any content.")
        # let the others raise by themselves
        # collect title if available
        tag_title = item.search_tag("pptx_slide_title")
        if tag_title:
            processed_data["item_ctx"]["tag_title"] = tag_title

        return {
            "guid": str(item.guid),
            **processed_data
        }

    @staticmethod
    def _split_data(data, limit):
        content = data["data"]
        split_list = []
        part_num = 0
        total_parts = len(content) // limit
        if len(content) % limit != 0:
            total_parts += 1
        for i in range(0, len(content), limit):
            data_split = copy.deepcopy(data)
            # split vertically and update context with metadata
            split_content = content[i:i + limit]
            if "num_lines" in data_split["item_ctx"]:
                data_split["item_ctx"]["num_lines"] = len(split_content)
            if "num_rows" in data_split["item_ctx"]:
                data_split["item_ctx"]["num_rows"] = len(split_content)
            data_split["data"] = split_content
            data_split["item_ctx"].update({
                "part_num": part_num,
                "total_parts": total_parts,
                "part_limit": limit
            })
            split_list.append(data_split)
            part_num += 1
        return split_list

    def _split_item_data(self, data, placeholder, params):
        item_type = data["type"]
        item_ctx = data["item_ctx"]
        split_list = []
        if item_type in ("string", "html"):
            num_lines = item_ctx["num_lines"]
            lines_per_slide = self.DefaultLimit.LINE
            # user setting if available
            if "lines" in params:
                try:
                    lines_per_slide = int(params["lines"])
                except ValueError:
                    raise Exception(f'The placeholder parameter "lines" is invalid')
            # cap at the lines available
            lines_per_slide = min(lines_per_slide, num_lines)
            data["item_ctx"]["num_lines"] = lines_per_slide
            # split item
            split_list = self._split_data(data, lines_per_slide)
        elif item_type in ("table", "tree"):
            # DEFAULT SPLITTING
            num_rows = item_ctx["num_rows"]
            num_cols = item_ctx["num_cols"]
            rows_per_slide = self.DefaultLimit.ROW
            cols_per_slide = self.DefaultLimit.COL
            # if user specified row/col count is available.
            if isinstance(placeholder, TablePlaceholder):
                if "rows" in params:
                    try:
                        rows_per_slide = int(params["rows"])
                    except ValueError:
                        raise Exception(f'The placeholder parameter "rows" is invalid')
                if "cols" in params:
                    try:
                        cols_per_slide = int(params["cols"])
                    except ValueError:
                        raise Exception(f'The placeholder parameter "cols" is invalid')
            else:  # isinstance(placeholder, PlaceholderGraphicFrame)
                table = placeholder.table
                rows_per_slide = len(table.rows)
                cols_per_slide = len(table.columns)
            # cap at the rows, cols available
            rows_per_slide = min(rows_per_slide, num_rows)
            cols_per_slide = min(cols_per_slide, num_cols)
            data["item_ctx"]["num_rows"] = rows_per_slide
            data["item_ctx"]["num_cols"] = cols_per_slide
            # split item
            split_list = self._split_data(data, rows_per_slide)
        return split_list

    @staticmethod
    def _create_item_from_string(contents, i_type):
        from data.models import Item
        item = Item()
        item.guid = uuid.uuid1()
        item.payloaddata = pickle.dumps(contents.strip(), protocol=0)
        item.type = i_type
        return item

    def _get_tag_title(self, items, ctx):
        if items is None:
            return
        # rebuild with tags
        for data in items:
            item_ctx = data["item_ctx"]
            tag_title = item_ctx.get("tag_title")
            if tag_title:
                tag_title = tag_title.strip()
                out_type = "html" if HTMLParser().is_html(tag_title) else "string"
                title_item = self._create_item_from_string(tag_title, out_type)
                return self._process_item(title_item, ctx)

    def _get_sec_shape_info(self, sec_shapes, ctx, items=None):
        if items is not None:
            items = list(items)
        processed_info = {}
        # process title first because breadcrumbs are dependent on it
        title_info = sec_shapes.get("title", None)
        if title_info:
            shape = title_info["shape"]
            shape_id = f"{shape.shape_id}"
            query_items = title_info["items"]
            if query_items:
                item = query_items[-1]
                # match item with ph and raise before processing
                self._verify_type_and_raise(shape, item)
                title_data = self._process_item(item, {**ctx, "title": True})
                title = title_data["item_ctx"]["title"]
                # save titles for all slides except the TOC slide
                if not ctx.get("toc", False):
                    slide = self.pptx_mgr.get_current_slide(shape)
                    self.title_info[f"{slide.slide_id}"] = title
                processed_info[shape_id] = self._get_tag_title(items, ctx) or title_data
        # other shapes
        for shape_name, shape_info in sec_shapes.items():
            shape = shape_info["shape"]
            shape_id = f"{shape.shape_id}"
            query_items = shape_info["items"]
            item = None
            shape_ctx = {}
            if shape_name == self.Param.TOC_LINK:
                # TOC slide must come before other slides that have a TOC link.
                if self.toc_slide_id is None:
                    continue
                item = query_items[-1] if query_items else \
                    self._create_item_from_string(
                        self.TOC_LINK_DEFAULT_TEXT,
                        "string"
                    )
                shape_ctx = {
                    "extra_info": {
                        0: {
                            "link_type": self.HyperlinkType.SLIDE,
                            "link_target": self.toc_slide_id,
                        }
                    }
                }
            elif shape_name == self.Param.COMMENTS:
                if "comments" in ctx:
                    item = self._create_item_from_string(ctx["comments"], "html")
            elif shape_name == self.Param.BREADCRUMB:
                slide = self.pptx_mgr.get_current_slide(shape)
                slide_id = f"{slide.slide_id}"
                breadcrumb = ""
                extra_info = {}
                if slide_id in self.hierarchy_info:
                    ancestors = self.hierarchy_info[slide_id]["ancestors"]
                    if ancestors:
                        curr_idx = 0
                        for sl_id in ancestors:
                            breadcrumb += self.title_info.get(sl_id, "")
                            extra_info[curr_idx] = {
                                "link_type": "slide",
                                "link_target": sl_id,
                            }
                            breadcrumb += "\n / \n"
                            # +1 for the crumb, +1 for the '/'
                            curr_idx += 2
                # add current title
                breadcrumb += self.title_info.get(slide_id, "")
                # add part numbers
                if items:
                    part_num = 0
                    total_parts = 0
                    # if there are multiple items, consider them as one
                    # entity and get max value for page count
                    for item_data in items:
                        if not item_data:
                            continue
                        item_ctx = item_data["item_ctx"]
                        part_num = max(item_ctx.get('part_num', 0), part_num)
                        total_parts = max(item_ctx.get('total_parts', 1), total_parts)
                    if total_parts > 1:
                        breadcrumb += f" ({part_num + 1} of {total_parts})"
                if breadcrumb:
                    item = self._create_item_from_string(breadcrumb, "string")
                    shape_ctx = {"extra_info": extra_info, "as_run": True}

            if item is not None:
                self._verify_type_and_raise(shape, item)
                processed_info[shape_id] = self._process_item(item, ctx | shape_ctx)

        return processed_info

    def _process_slide(self, slide, items, ctx, as_toc=False):
        from data.models import Item
        slide_id = f"{slide.slide_id}"
        # shape_info is a map of shapes and the item to be inserted into each shape, possibly split
        # into parts if an item exceeds a slide's limit.
        shape_info = {}
        # similar to shape_info but for shapes that have only one item or part to insert.
        # These must be repeated across slides that were duplicated from the current source slide.
        repeated_shape_info = {}
        # shapes that are not the actual content
        secondary_shapes = {}
        # NOTE: A slide can have multiple shapes.
        # In powerpoint terms, a 'placeholder' is any
        # shape that does not have any content inserted.
        # We accept all types of shapes - with or without
        # content inserted BUT, they must belong in the allowed
        # list. Due to this reason, the terms 'placeholder' and
        # 'shape' are used interchangeable in the code.
        for shape in slide.shapes:
            # only accept select shapes
            if shape.placeholder_format.type not in self.ALLOWED_PLACEHOLDERS:
                continue
            selection_label = unescape(shape.name)
            if not selection_label.startswith(self.PARAM_PREFIX):
                continue
            shape_id = f"{shape.shape_id}"
            # build params
            params_str = selection_label.replace(self.PARAM_PREFIX, "").strip()
            params = self._parse_params(params_str)
            # shape-local context
            shape_ctx = {}
            # find the items to insert
            items_to_insert = []
            secondary = False
            # - A toc shape is a text placeholder where
            # titles of other slides are inserted as a
            # list of lines hyperlinking the corresponding
            # slide.
            # - The presence of this shape tells us the slide
            # must be categorized as a toc slide.
            if self.Param.TOC in params:
                # Process TOC slides at the end
                if not as_toc:
                    return True
                # Prepare a string item on the fly
                # using a list of all the slide titles.
                str_content = ""
                extra_info = {}
                idx = 0
                for sl_id, title in self.title_info.items():
                    if sl_id in self.toc_exclude_set:
                        continue
                    extra_info[idx] = {
                        "link_type": self.HyperlinkType.SLIDE,
                        "link_target": sl_id,
                    }
                    if sl_id in self.hierarchy_info:
                        extra_info[idx]["depth"] = self.hierarchy_info[sl_id]["depth"]
                    str_content += title
                    if idx != len(self.title_info) - 1:
                        str_content += "\n"
                    idx += 1

                items_to_insert = [self._create_item_from_string(str_content, "string")]
                shape_ctx["extra_info"] = extra_info
            else:
                # check query for other cases.
                query = params.pop("query", None)
                # go ahead and query now
                if query is not None:
                    items_to_insert = list(Item.find(ctx['request'], query=query, queryset=items))

                # handle certain shapes later
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE or self.Param.HTML_HEADER in params:  # titles
                    item = None
                    if items_to_insert:
                        item = items_to_insert[-1]
                    else:
                        if self.Param.HTML_HEADER in params and "html_header" in ctx:
                            item = self._create_item_from_string(ctx["html_header"], "html")
                    if item is not None:
                        secondary = True
                        secondary_shapes["title"] = {"shape": shape, "items": [item]}
                else:  # others
                    for param in (
                            self.Param.TOC_LINK,
                            self.Param.BREADCRUMB,
                            self.Param.COMMENTS
                    ):
                        if param in params:
                            secondary = True
                            secondary_shapes[str(param)] = {"shape": shape, "items": items_to_insert}
                            break

            # skip
            if secondary or not items_to_insert:
                continue

            # add content
            content = []
            item_type = None
            for item in items_to_insert:
                # match item with ph and raise before processing
                self._verify_type_and_raise(shape, item)
                # process
                data = self._process_item(item, {**ctx, **shape_ctx})
                # first item in the valid list decides the type to be inserted for this shape.
                if item_type is None:
                    item_type = data["type"]
                else:
                    if data["type"] != item_type:
                        raise Exception(f"Items to be inserted in the placeholder with label '{selection_label}'"
                                        f" must be of the same type.")
                # only selected types must be split
                if item_type not in self.ITEM_SPLIT_ALLOWED_TYPES:
                    # just one part
                    item_parts = [data]
                else:
                    # split item based on ctx/placeholder
                    # can also return just one part, meaning no split required
                    item_parts = self._split_item_data(data, shape, params)
                # add parts
                for part in item_parts:
                    content.append(part)
            if len(content) == 1:
                # If there is only item/part, add to the repeated info
                # which would then be copied to every duplicated slide
                repeated_shape_info[f"{shape_id}"] = content[0]
            else:
                # if an item is multipart or if more than one item.
                shape_info[f"{shape_id}"] = content
            # Escaping characters is necessary or else XML writing will fail
            # when the presentation is saved.
            shape.name = escape(selection_label)

        # slide_info is where we store the mapping between a placeholder and the item
        # that goes in there. It is of the form:
        # {
        #   "shape1": {..item1_data..},
        #   "shape2": {..item2_data..}
        # }
        # Check fill_slide() on how this is used.
        slide_info = {}

        if shape_info or repeated_shape_info:
            # - shape_info is of the form:
            # shape_info = {
            # 'SHAPE-2': [I1P1, I1P2, I1P3],
            # 'SHAPE-4': [I3P1],
            # 'SHAPE-6': [I4P1, I4P2]
            # }
            # where I1P4 -> Item1 Part4 which is a dict
            # containing parts of the item data based on per-slide limits
            #
            # - repeated_shape_info is of the form:
            # repeated_shape_info = {
            # 'SHAPE-1': I1P1,
            # 'SHAPE-3': I3P1,
            # 'SHAPE-5': I4P1
            # }
            # Note that each shape has only one item to insert. This is to
            # make sure shapes/placeholders in duplicated slides are filled
            # by the same item.
            if shape_info:
                # Now: Divide items among slides by zipping across shape info.
                # output looks like: [(I1P1, I3P1, I4P1), (I1P2, None, I4P2), (I1P3, None, None)]
                slide_item_tuples = list(zip_longest(*shape_info.values()))
                # get shape ids
                shape_ids = shape_info.keys()
                duplicate_count = -1
                # build slide info map
                # len(slide_item_tuples) is the number of slides required.
                # So each iteration is for one slide
                for item_tuple in slide_item_tuples:
                    duplicate_count += 1
                    # eg: for idx 0 -> { 'SHAPE-2': I1P1, 'SHAPE-4': I3P1, 'SHAPE-6': None }
                    content = dict(zip(shape_ids, item_tuple))
                    # include repeated shape data as well
                    content.update(repeated_shape_info)
                    # MUST deepcopy to avoid editing the same object.
                    info = copy.deepcopy(slide_info)
                    secondary_shape_info = self._get_sec_shape_info(secondary_shapes,
                                                                    {**ctx, "toc": as_toc},
                                                                    items=content.values())
                    # add a temporary slide id for duplicates
                    slide_key = f"{slide_id}-{duplicate_count}" if duplicate_count > 0 else slide_id
                    # add
                    self.ppt_info[slide_key] = info | secondary_shape_info | content
                # record the number of times to duplicate a slide.
                if duplicate_count > 0:
                    self.duplicate_count_info[slide_id] = duplicate_count
            else:  # if only one item/ one part to insert
                info = copy.deepcopy(slide_info)
                secondary_shape_info = self._get_sec_shape_info(secondary_shapes,
                                                                {**ctx, "toc": as_toc},
                                                                items=repeated_shape_info.values())
                self.ppt_info[slide_id] = info | secondary_shape_info | repeated_shape_info
        else:
            # write also if there's just the title
            secondary_shape_info = self._get_sec_shape_info(secondary_shapes, {**ctx, "toc": as_toc})
            slide_info.update(secondary_shape_info)
            self.ppt_info[slide_id] = slide_info

        return False

    def _reproduce_slide(self, idx, slide, count):
        # Reproduces a slide "count" times
        curr_idx = idx
        # duplicate
        for i in range(count):
            # duplicate at next index
            curr_idx = curr_idx + 1
            try:
                copied_slide = self.pptx_mgr.duplicate_slide(slide, curr_idx)
                # update the fake id with the actual id after slide creation
                duplicate_key = f"{slide.slide_id}-{i + 1}"
                if duplicate_key in self.ppt_info:
                    slide_info = self.ppt_info.pop(duplicate_key)
                    self.ppt_info[f"{copied_slide.slide_id}"] = slide_info
            except Exception as e:
                raise Exception(f'Error duplicating slides :: {str(e)}')

    def _fill_slide(self, slide):
        # Fill the slide with processed data
        slide_id = slide.slide_id
        slide_info = self.ppt_info.get(f"{slide_id}")
        if slide_info:
            for shape in slide.shapes:
                shape_id = f"{shape.shape_id}"
                if shape_id in slide_info:
                    data = slide_info[shape_id]
                    if data:
                        self._add_item(shape, data)

    def generate_report(self, parent, items, ctx):
        from data.models import Item
        # check if we want to keep all slides
        use_all_slides = parent.get_default(ctx, "use_all_slides", 1, force_int=True) == 1
        # this ignores all children templates and their filters
        if use_all_slides:
            # convert items to a queryset for later filtering
            child_qs = Item.filtered_objects.filter(pk__in=[item.guid for item in items])
            for slide in self.pptx_mgr.slides:
                is_toc = self._process_slide(slide, child_qs, ctx)
                # toc? save the id once
                if self.toc_slide_id is None and is_toc:
                    self.toc_slide_id = f"{slide.slide_id}"
            # Process TOC separately in the end.
            if self.toc_slide_id is not None:
                toc_slide = self.pptx_mgr.get_slide(self.toc_slide_id)
                self._process_slide(toc_slide, child_qs, ctx, as_toc=True)
        else:  # pick slides
            children = parent.get_child_layouts(recursive=True)
            if not children:
                raise Exception(f"use_all_slides is set to 0 but the PPTX Layout template"
                                f" '{parent.name}' does not have any children")
            # process layouts
            curr_slide_count = len(self.pptx_mgr.slides)
            final_slide_indices = []
            toc_items = None
            toc_ctx = None
            template_slide_map = {}
            for idx, child in enumerate(children):
                child_ctx, child_items = child.setup_template(ctx, items)
                if child_items is None:
                    continue
                source_slide = child.get_default(child_ctx, 'source_slide', default=idx + 1, force_int=True)
                source_slide_idx = source_slide - 1
                if source_slide_idx in final_slide_indices:
                    # todo
                    raise Exception(
                        f"A slide can only be picked once from the input pptx."
                        f" 'source_slide' is set to '{source_slide}' more than once.")
                if source_slide > curr_slide_count:
                    raise Exception(
                        f"The input pptx only has {curr_slide_count} slides but the value of"
                        f" 'source_slide' for template '{child.name}' is {source_slide}.")
                # collect
                final_slide_indices.append(source_slide_idx)
                # add in html header if available.
                html_header = child.parse_HTML(child_ctx)
                if html_header:
                    child_ctx["html_header"] = html_header
                # add in comments if available.
                comments = child.parse_comments(child_ctx)
                if comments:
                    child_ctx["comments"] = comments
                # get slide and process
                slide = self.pptx_mgr.slides[source_slide_idx]
                slide_id = f"{slide.slide_id}"
                template_slide_map[child.template.guid] = slide_id
                # use template hierarchy
                depth = 0
                ancestors = []
                while child.parent is not None and child.parent is not parent:
                    # parent must be something other than the PPTX root layout
                    depth += 1
                    # parents are processed first so their slide ids are already available.
                    ancestors.insert(0, template_slide_map[child.parent.template.guid])
                    child = child.parent
                # add
                self.hierarchy_info[slide_id] = {
                    "depth": depth,
                    "ancestors": ancestors
                }
                # convert items to a queryset for later filtering
                child_qs = Item.filtered_objects.filter(pk__in=[item.guid for item in child_items])
                # process
                is_toc = self._process_slide(slide, child_qs, child_ctx)
                # toc
                exclude_toc = child.get_default(child_ctx, "exclude_from_toc", default=0, force_int=True) == 1
                if exclude_toc:
                    self.toc_exclude_set.add(slide_id)
                if self.toc_slide_id is None and is_toc:
                    # save the toc slide for later
                    self.toc_slide_id = slide_id
                    toc_items = child_qs
                    toc_ctx = child_ctx
            # Process TOC separately in the end.
            if self.toc_slide_id is not None:
                toc_slide = self.pptx_mgr.get_slide(self.toc_slide_id)
                self._process_slide(toc_slide, toc_items, toc_ctx, as_toc=True)
            # delete unwanted slides
            for idx, slide in enumerate(self.pptx_mgr.slides):
                if idx not in final_slide_indices:
                    self.pptx_mgr.delete_slide(slide)
            # necessary step after deletion to
            # rebuild internal relationships
            self.pptx_mgr.save_presentation()
        # rebuild
        if self.duplicate_count_info:
            total_duplicates = 0
            for slide_idx, slide in enumerate(self.pptx_mgr.slides):
                # duplicate
                slide_id = f"{slide.slide_id}"
                if slide_id in self.duplicate_count_info:
                    duplication_count = self.duplicate_count_info[slide_id]
                    new_idx = slide_idx + total_duplicates
                    self._reproduce_slide(new_idx, slide, duplication_count)
                    total_duplicates += duplication_count
        # populate
        for slide in self.pptx_mgr.slides:
            self._fill_slide(slide)
        # save
        self.pptx_mgr.save_presentation()

    def get_report_content(self):
        return self.pptx_mgr.get_presentation_stream()
