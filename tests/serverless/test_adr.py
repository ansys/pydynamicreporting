from pathlib import Path
from random import random as r
import uuid

from django.core.exceptions import ImproperlyConfigured
import numpy as np
import pytest

from ansys.dynamicreporting.core.exceptions import (
    ADRException,
    ImproperlyConfiguredError,
    InvalidAnsysPath,
    InvalidPath,
)
from ansys.dynamicreporting.core.serverless import ADR


@pytest.mark.ado_test
def test_create_no_setup():
    from ansys.dynamicreporting.core.serverless import Session

    with pytest.raises(RuntimeError):
        Session.create()


@pytest.mark.ado_test
def test_get_instance_error():
    with pytest.raises(RuntimeError):
        ADR.get_instance()


@pytest.mark.ado_test
def test_get_instance_no_setup():
    from ansys.dynamicreporting.core.constants import DOCKER_DEV_REPO_URL

    adr = ADR(  # noqa: F841
        ansys_installation="docker",
        docker_image=DOCKER_DEV_REPO_URL,
        media_url="/media1/",
        static_url="/static2/",
        in_memory=True,
    )
    assert ADR.get_instance() is adr


@pytest.mark.ado_test
def test_ensure_setup_error_no_setup():
    from ansys.dynamicreporting.core.constants import DOCKER_DEV_REPO_URL

    adr = ADR(  # noqa: F841
        ansys_installation="docker",
        docker_image=DOCKER_DEV_REPO_URL,
        media_url="/media1/",
        static_url="/static2/",
        in_memory=True,
    )

    with pytest.raises(RuntimeError):
        ADR.ensure_setup()


@pytest.mark.ado_test
def test_get_database_config_before_setup():
    assert ADR.get_database_config() is None


@pytest.mark.ado_test
def test_get_database_config_before_setup_raise():
    with pytest.raises(ImproperlyConfiguredError):
        ADR.get_database_config(raise_exception=True)


@pytest.mark.ado_test
def test_is_setup_before_setup():
    assert not ADR.get_instance().is_setup


@pytest.mark.ado_test
def test_get_instance(adr_serverless):
    assert ADR.get_instance() is adr_serverless


@pytest.mark.ado_test
def test_is_setup_after_setup(adr_serverless):
    assert adr_serverless.is_setup


@pytest.mark.ado_test
def test_has_static_files(adr_serverless):
    static_dir = adr_serverless.static_directory
    assert (Path(static_dir) / "admin" / "css" / "base.css").exists()


@pytest.mark.ado_test
def test_get_database_config_after_setup(adr_serverless):
    assert "default" in adr_serverless.get_database_config()


@pytest.mark.ado_test
def test_setup_after_setup(adr_serverless):
    with pytest.raises(RuntimeError):
        adr_serverless.setup(collect_static=True)


@pytest.mark.ado_test
def test_get_instance_session():
    from ansys.dynamicreporting.core.serverless import Session

    adr = ADR.get_instance()
    assert adr.session is not None and isinstance(adr.session, Session) and adr.session_guid


@pytest.mark.ado_test
def test_get_instance_dataset():
    from ansys.dynamicreporting.core.serverless import Dataset

    adr = ADR.get_instance()
    assert adr.dataset is not None and isinstance(adr.dataset, Dataset) and adr.dataset.guid


@pytest.mark.ado_test
def test_init_twice(adr_serverless):
    # return the same instance
    from ansys.dynamicreporting.core.constants import DOCKER_DEV_REPO_URL

    adr = ADR(
        ansys_installation="docker",
        docker_image=DOCKER_DEV_REPO_URL,
        db_directory=adr_serverless.db_directory,
        static_directory=adr_serverless.static_directory,
        media_url="/media1/",
        static_url="/static2/",
    )
    assert adr is adr_serverless


@pytest.mark.ado_test
def test_set_default_session(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Session

    session = Session.create(application="serverless test default sesh", tags="dp=dp227")
    adr_serverless.set_default_session(session)
    assert adr_serverless.session_guid == session.guid


@pytest.mark.ado_test
def test_set_default_session_no_session(adr_serverless):
    with pytest.raises(TypeError, match="Must be an instance of type 'Session'"):
        adr_serverless.set_default_session(None)


@pytest.mark.ado_test
def test_set_default_dataset(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Dataset

    dataset = Dataset.create(filename="serverless test default dataset", tags="dp=dp227")
    adr_serverless.set_default_dataset(dataset)
    assert adr_serverless.dataset.guid == dataset.guid


@pytest.mark.ado_test
def test_set_default_dataset_no_dataset(adr_serverless):
    with pytest.raises(TypeError, match="Must be an instance of type 'Dataset'"):
        adr_serverless.set_default_dataset(None)


@pytest.mark.ado_test
def test_create_html(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML

    # HTML item
    item_html = (
        "<h1>Heading 1</h1>"
        "<h2>Heading 2</h2>"
        "<h3>Heading 3</h3>"
        "<h4>Heading 4</h4>"
        "<h5>Heading 5</h5>"
        "Two breaks below"
        "<br><br />"
        "<h6>Heading 6 (& one break below)</h6>"
        "<br>"
        "The end"
    )
    intro_html = adr_serverless.create_item(
        HTML,
        name="intro_html",
        content=item_html,
        source="sls-test",
    )
    assert HTML.get(name="intro_html").guid == intro_html.guid


@pytest.mark.ado_test
def test_create_item_type_failure(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Template

    with pytest.raises(TypeError, match="Template is not a subclass of Item"):
        adr_serverless.create_item(Template, name="intro_html", content="lololol")


@pytest.mark.ado_test
def test_create_item_kwarg_failure(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML

    with pytest.raises(TypeError, match="got an unexpected keyword argument 'foo'"):
        adr_serverless.create_item(
            HTML,
            name="intro_html",
            content="<h1>Heading 1</h1>",
            foo="bar",
        )


@pytest.mark.ado_test
def test_create_item_empty_kwarg_failure(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML

    with pytest.raises(ADRException, match="At least one keyword argument must be provided"):
        adr_serverless.create_item(HTML)


@pytest.mark.ado_test
def test_edit_html(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML

    adr_serverless.create_item(
        HTML,
        name="test_edit_html",
        content="<h1>Heading 1</h1>",
    )
    intro_html = HTML.get(name="test_edit_html")
    intro_html.content = "<h2>Heading 2</h2>" "<br>"
    intro_html.save()

    assert "h1" not in HTML.get(guid=intro_html.guid).content


@pytest.mark.ado_test
def test_create_html_bad_content(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML

    with pytest.raises(ValueError):  # must be valid HTML
        adr_serverless.create_item(HTML, name="empty_html", content="lololol")


@pytest.mark.ado_test
def test_create_string(adr_serverless):
    from ansys.dynamicreporting.core.serverless import String

    # string
    item_text = "This section describes the settings for the simulation: initial conditions, solver settings, and such."
    intro_text = adr_serverless.create_item(
        String,
        name="intro_text",
        content=item_text,
        tags="dp=dp227 section=intro",
        source="sls-test",
    )

    assert String.get(name="intro_text").guid == intro_text.guid


@pytest.mark.ado_test
def test_create_anim(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Animation

    # anim
    anim = adr_serverless.create_item(
        Animation,
        name="intro_anim",
        content=str(Path(__file__).parent / "test_data" / "movie.mp4"),
        tags="dp=dp227 section=data",
        source="sls-test",
        sequence=1,
    )

    assert Animation.get(name="intro_anim").guid == anim.guid


@pytest.mark.ado_test
def test_create_anim_no_file(adr_serverless):  # file does not exist
    from ansys.dynamicreporting.core.serverless import Animation

    with pytest.raises(ValueError):
        adr_serverless.create_item(
            Animation,
            name="bad_anim",
            content=str(Path(__file__).parent / "test_data" / "lolol.mp4"),
        )


@pytest.mark.ado_test
def test_create_anim_bad_format(adr_serverless):  # file is not a video
    from ansys.dynamicreporting.core.serverless import Animation

    with pytest.raises(ValueError):
        adr_serverless.create_item(
            Animation,
            name="bad_anim",
            content=str(Path(__file__).parent / "test_data" / "scene.avz"),
        )


@pytest.mark.ado_test
def test_create_file(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File

    # file
    file = adr_serverless.create_item(
        File,
        name="intro_file",
        content=str(Path(__file__).parent / "test_data" / "input.pptx"),
        tags="dp=dp227 section=data",
        source="sls-test",
        sequence=1,
    )

    assert File.get(name="intro_file").guid == file.guid


@pytest.mark.ado_test
def test_create_file_csf(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File

    # csf
    intro_csf = adr_serverless.create_item(
        File,
        name="intro_csf",
        content=str(Path(__file__).parent / "test_data" / "scene1.csf"),
        tags="dp=dp227 section=data",
        source="sls-test",
        sequence=1,
    )

    assert File.get(name="intro_csf").guid == intro_csf.guid


@pytest.mark.ado_test
def test_create_file_ens(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File

    # ens
    intro_ens = adr_serverless.create_item(
        File,
        name="intro_ens",
        content=str(Path(__file__).parent / "test_data" / "scene2.ens"),
        tags="dp=dp227 section=data",
        source="sls-test",
        sequence=1,
    )

    assert File.get(name="intro_ens").guid == intro_ens.guid


@pytest.mark.ado_test
def test_create_file_evsn(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File

    # evsn
    intro_evsn = adr_serverless.create_item(
        File,
        name="intro_evsn",
        content=str(Path(__file__).parent / "test_data" / "scenario.evsn"),
        tags="dp=dp227 section=data",
        source="sls-test",
        sequence=1,
    )

    assert File.get(name="intro_evsn").guid == intro_evsn.guid


@pytest.mark.ado_test
def test_create_image(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Image

    # image
    intro_image = adr_serverless.create_item(
        Image,
        name="intro_image",
        content=str(Path(__file__).parent / "test_data" / "nexus_logo.png"),
        tags="dp=dp227 section=data",
        source="sls-test",
    )

    assert Image.get(name="intro_image").guid == intro_image.guid


@pytest.mark.ado_test
def test_create_enhanced_image(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Image

    # enhanced image
    intro_enhanced_image = adr_serverless.create_item(
        Image,
        name="intro_enhanced_image",
        content=str(Path(__file__).parent / "test_data" / "case.tif"),
        tags="dp=dp227 section=data",
        source="sls-test",
    )

    assert Image.get(name="intro_enhanced_image").guid == intro_enhanced_image.guid


@pytest.mark.ado_test
def test_create_scene(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Scene

    # scene
    scene = adr_serverless.create_item(
        Scene,
        name="intro_scene",
        content=str(Path(__file__).parent / "test_data" / "scene.avz"),
        session=adr_serverless.session,  # random test
        dataset=adr_serverless.dataset,
        tags="dp=dp227 section=data",
        source="sls-test",
    )

    assert Scene.get(name="intro_scene").guid == scene.guid


@pytest.mark.ado_test
def test_create_table(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Table

    # table
    ics = []
    ips = []
    zet = []
    for i in range(30):
        ics.append(i / 5.0)
        ips.append(np.sin((i + 6 * 0) * np.pi / 10.0) + r() * 0.1)
        zet.append(np.cos((i + 6 * 0) * np.pi / 10.0) + r() * 0.1)

    intro_table = adr_serverless.create_item(
        Table,
        name="intro_table",
        content=np.array([ics, ips, zet], dtype="|S20"),
        tags="dp=0 type=hex8",
        source="sls-test",
    )

    intro_table.labels_row = ["X", "Sin", "Cos"]
    intro_table.set_tags("dp=dp227 section=data")
    intro_table.plot = "line"
    intro_table.xaxis = "X"
    intro_table.yaxis = ["Sin", "Cos"]
    intro_table.xaxis_format = "floatdot0"
    intro_table.yaxis_format = "floatdot1"
    intro_table.ytitle = "Values"
    intro_table.xtitle = "X"

    intro_table.save()

    assert Table.get(name="intro_table").guid == intro_table.guid


@pytest.mark.ado_test
def test_create_tree(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Tree

    # tree
    tree_content = [
        {"key": "root", "name": "Solver", "value": "My Solver"},
        {"key": "root", "name": "Number cells", "value": 10e6},
        {"key": "root", "name": "Mesh Size", "value": "1.0 mm^3"},
        {"key": "root", "name": "Mesh Type", "value": "Hex8"},
    ]
    # alternative way of creation
    tree = adr_serverless.create_item(
        Tree,
        name="intro_tree",
        content=tree_content,
        tags="dp=dp227 section=data",
        session=adr_serverless.session,
        dataset=adr_serverless.dataset,
        source="sls-test",
    )

    assert Tree.get(name="intro_tree").guid == tree.guid


@pytest.mark.parametrize(
    "backup_kwargs",
    [
        {"compress": True},
        {"ignore_primary_keys": True},
    ],
)
@pytest.mark.ado_test
def test_backup_success(adr_serverless, tmp_path, backup_kwargs):
    adr_serverless.backup_database(output_directory=str(tmp_path), **backup_kwargs)
    json_files = list(tmp_path.glob("*.gz" if "compress" in backup_kwargs else "*.json"))
    assert any(f.name.startswith("backup_") for f in json_files)


@pytest.mark.ado_test
def test_backup_in_memory_disallowed(adr_serverless, tmp_path, monkeypatch):
    monkeypatch.setattr(adr_serverless, "_in_memory", True)
    with pytest.raises(ADRException, match="Backup is not available in in-memory mode."):
        adr_serverless.backup_database(output_directory=str(tmp_path))


@pytest.mark.ado_test
def test_backup_invalid_output_directory(adr_serverless, tmp_path):
    # test path object and file at the same time
    random_file = tmp_path / "not_created_yet.txt"
    random_file.touch(exist_ok=True)
    with pytest.raises(InvalidPath, match="not a valid directory"):
        adr_serverless.backup_database(output_directory=random_file)


@pytest.mark.ado_test
def test_backup_invalid_database(adr_serverless, tmp_path):
    with pytest.raises(ADRException, match="must be configured first"):
        adr_serverless.backup_database(output_directory=tmp_path, database="not_a_database")


@pytest.mark.ado_test
def test_restore_invalid_database(adr_serverless, tmp_path):
    with pytest.raises(ADRException, match="must be configured first"):
        base_dir = Path(__file__).parent / "test_data"
        adr_serverless.restore_database(str(base_dir / "restoreme.json"), database="not_a_database")


@pytest.mark.ado_test
def test_restore_invalid_file_path(adr_serverless, tmp_path):
    with pytest.raises(InvalidPath, match="not a valid file"):
        adr_serverless.restore_database(tmp_path)


@pytest.mark.ado_test
def test_restore_backup_success(adr_serverless):
    base_dir = Path(__file__).parent / "test_data"
    # should restore without error
    adr_serverless.restore_database(str(base_dir / "restoreme.json"))


@pytest.mark.ado_test
def test_backup_django_command_failure(adr_serverless, tmp_path, monkeypatch):
    from ansys.dynamicreporting.core.serverless import adr as adr_module

    def fake_call_command(*args, **kwargs):
        raise Exception("backup error")

    monkeypatch.setattr(adr_module, "call_command", fake_call_command)
    with pytest.raises(ADRException):
        adr_serverless.backup_database(output_directory=str(tmp_path))


@pytest.mark.ado_test
def test_restore_django_command_failure(adr_serverless, tmp_path, monkeypatch):
    from ansys.dynamicreporting.core.serverless import adr as adr_module

    def fake_call_command(*args, **kwargs):
        raise Exception("restore error")

    monkeypatch.setattr(adr_module, "call_command", fake_call_command)
    with pytest.raises(ADRException):
        json_file = tmp_path / "bad.json"
        json_file.write_text('[{"invalid": "bad"}]')
        adr_serverless.restore_database(str(json_file))


@pytest.mark.ado_test
def test_get_ansys_installation(adr_serverless):
    assert adr_serverless.ansys_installation is not None and isinstance(
        adr_serverless.ansys_installation, str
    )


@pytest.mark.ado_test
def test_get_ansys_version(adr_serverless):
    assert adr_serverless.ansys_version is not None and isinstance(
        adr_serverless.ansys_version, int
    )


@pytest.mark.ado_test
def test_get_media_directory(adr_serverless):
    assert adr_serverless.media_directory is not None and isinstance(
        adr_serverless.media_directory, str
    )


@pytest.mark.ado_test
def test_create_demo_report(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout

    top_parent = adr_serverless.create_template(
        BasicLayout, name="Serverless Simulation Report", parent=None, tags="dp=dp227"
    )
    top_parent.params = '{"HTML": "<h1>Serverless Simulation Report</h1>"}'
    top_parent.set_filter("A|i_tags|cont|dp=dp227;")
    top_parent.save()

    from ansys.dynamicreporting.core.serverless import TOCLayout

    toc_layout = adr_serverless.create_template(
        TOCLayout, name="TOC", parent=top_parent, tags="dp=dp227"
    )
    toc_layout.params = '{"TOCitems": 1, "HTML": "<h2>Table of Content</h2>"}'
    toc_layout.set_filter("A|i_name|eq|__NonexistentName__;")
    toc_layout.save()

    from ansys.dynamicreporting.core.serverless import PanelLayout

    intro_panel = adr_serverless.create_template(
        PanelLayout, name="Introduction", parent=top_parent, tags="dp=dp227"
    )
    intro_panel.params = '{"HTML": "<h2>Introduction</h2>", "properties": {"TOCItem": "1"}}'
    intro_panel.set_filter("A|i_tags|cont|section=intro;")
    intro_panel.save()

    from ansys.dynamicreporting.core.serverless import PanelLayout

    # alternate way of creation
    results_panel = PanelLayout.create(name="Results", parent=top_parent, tags="dp=dp227")
    results_panel.params = (
        '{"HTML": "<h2>Results</h2>\\nYour simulation results.", "properties": {"TOCItem": "1"}}'
    )
    results_panel.set_filter("A|i_tags|cont|section=data;")
    results_panel.save()
    # create_template() does this for you
    top_parent.children.append(results_panel)
    top_parent.save()

    from ansys.dynamicreporting.core.serverless import Template

    test_parent = Template.get(name="Serverless Simulation Report")
    assert test_parent.guid == top_parent.guid
    for child in test_parent.children:
        assert child.guid in top_parent.children_order


@pytest.mark.ado_test
def test_create_template_toc(adr_serverless):
    from ansys.dynamicreporting.core.serverless import TOCLayout

    toc_layout = adr_serverless.create_template(
        TOCLayout, name="test_create_template_toc", parent=None
    )
    toc_layout.set_filter("A|i_name|eq|__NonexistentName__;")
    toc_layout.save()

    assert TOCLayout.get(guid=toc_layout.guid).guid == toc_layout.guid


@pytest.mark.ado_test
def test_create_template_type_failure(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File

    with pytest.raises(TypeError, match="File is not a subclass of Template"):
        adr_serverless.create_template(File, name="test_create_template_toc", parent=None)


@pytest.mark.ado_test
def test_create_template_kwarg_failure(adr_serverless):
    from ansys.dynamicreporting.core.serverless import TOCLayout

    with pytest.raises(TypeError, match="got an unexpected keyword argument 'foo'"):
        adr_serverless.create_template(
            TOCLayout,
            name="test_create_template_toc",
            parent=None,
            foo="bar",
        )


@pytest.mark.ado_test
def test_create_template_kwarg_empty_failure(adr_serverless):
    from ansys.dynamicreporting.core.serverless import TOCLayout

    with pytest.raises(ADRException, match="At least one keyword argument must be provided"):
        adr_serverless.create_template(TOCLayout)


@pytest.mark.ado_test
def test_get_report_basic(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout

    top_parent = adr_serverless.create_template(BasicLayout, name="test_get_report", parent=None)
    top_parent.set_filter("A|i_name|eq|__NonexistentName__;")
    top_parent.save()

    report = adr_serverless.get_report(guid=top_parent.guid)
    assert report.guid == top_parent.guid


@pytest.mark.ado_test
def test_get_report_empty_kwarg_failure(adr_serverless):
    with pytest.raises(ADRException, match="At least one keyword argument must be provided"):
        adr_serverless.get_report()


@pytest.mark.ado_test
def test_get_report_invalid_guid(adr_serverless):
    with pytest.raises(ADRException, match="Report not found"):
        adr_serverless.get_report(guid=str(uuid.uuid4()))


@pytest.mark.ado_test
def test_get_reports(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout

    top_parent = adr_serverless.create_template(BasicLayout, name="test_get_reports", parent=None)
    top_parent.set_filter("A|i_name|eq|__NonexistentName__;")
    top_parent.save()

    reports = adr_serverless.get_reports()
    assert len(reports) > 0


@pytest.mark.ado_test
def test_get_reports_w_fields(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout

    top_parent = adr_serverless.create_template(
        BasicLayout, name="test_get_reports_w_fields", parent=None
    )
    top_parent.set_filter("A|i_name|eq|__NonexistentName__;")
    top_parent.save()

    reports = adr_serverless.get_reports(fields=["guid", "name"])
    assert len(reports) > 0
    assert all(
        [
            isinstance(uuid.UUID(report[0]), uuid.UUID) and isinstance(report[1], str)
            for report in reports
        ]
    )


@pytest.mark.ado_test
def test_get_list_reports(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout, TOCLayout

    top_parent = adr_serverless.create_template(
        BasicLayout, name="test_get_list_reports", parent=None
    )
    top_parent.set_filter("A|i_name|eq|__NonexistentName__;")
    top_parent.save()

    toc_layout = adr_serverless.create_template(
        TOCLayout, name="test_get_list_reports_toc", parent=top_parent
    )
    toc_layout.set_filter("A|i_name|eq|__NonexistentName__;")
    toc_layout.save()

    reports = adr_serverless.get_list_reports()
    assert len(reports) > 0


@pytest.mark.ado_test
def test_get_list_reports_w_r_type(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout

    top_parent = adr_serverless.create_template(
        BasicLayout, name="test_get_list_reports_w_r_type", parent=None
    )
    top_parent.set_filter("A|i_name|eq|__NonexistentName__;")
    top_parent.save()
    reports = adr_serverless.get_list_reports(r_type="name")
    assert len(reports) > 0
    assert all([isinstance(name, str) for name in reports])


@pytest.mark.ado_test
def test_get_list_reports_w_r_type_reports(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout, Template

    top_parent = adr_serverless.create_template(
        BasicLayout, name="test_get_list_reports_w_r_type_reports", parent=None
    )
    top_parent.set_filter("A|i_name|eq|__NonexistentName__;")
    top_parent.save()
    reports = adr_serverless.get_list_reports(r_type=None)
    assert len(reports) > 0
    assert all([isinstance(rep, Template) for rep in reports])


@pytest.mark.ado_test
def test_get_list_reports_wrong_type(adr_serverless):
    with pytest.raises(ADRException, match="r_type must be one of"):
        adr_serverless.get_list_reports(r_type="wrong_type")


@pytest.mark.ado_test
def test_query_items(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File, Item

    # file
    adr_serverless.create_item(
        File,
        name="query_test_file",
        content=str(Path(__file__).parent / "test_data" / "input.pptx"),
    )

    objs = adr_serverless.query(query_type=File, query="A|i_name|cont|query_test_file;")
    objs2 = adr_serverless.query(
        query_type=Item, query="A|i_type|cont|file;A|i_name|cont|query_test_file;"
    )
    assert len(objs) == 1 and len(objs2) == 1


@pytest.mark.ado_test
def test_query_items_wrong_type(adr_serverless):
    from ansys.dynamicreporting.core.serverless import ADR

    with pytest.raises(TypeError, match="is not a type of"):
        adr_serverless.query(query_type=ADR, query="A|i_name|cont|query_test_file;")


@pytest.mark.ado_test
def test_query_templates(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Template, TOCLayout

    toc_layout = adr_serverless.create_template(TOCLayout, name="test_query_templates")
    toc_layout.set_filter("A|i_name|eq|__NonexistentName__;")
    toc_layout.save()

    temps = adr_serverless.query(query_type=TOCLayout, query="A|t_name|eq|test_query_templates;")
    temps2 = adr_serverless.query(
        query_type=Template, query="A|t_types|cont|Layout:toc;A|t_name|eq|test_query_templates;"
    )

    assert len(temps) == 1 and len(temps2) == 1


@pytest.mark.ado_test
def test_query_no_templates(adr_serverless):
    from ansys.dynamicreporting.core.serverless import PPTXLayout, Template

    temps = adr_serverless.query(query_type=PPTXLayout, query="A|t_name|eq|pptx-select;")
    temps2 = adr_serverless.query(
        query_type=Template, query="A|t_types|cont|Layout:pptx;A|t_name|eq|pptx-select;"
    )

    assert not temps and not temps2


@pytest.mark.ado_test
def test_create_objects(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML, String

    kwargs = {
        "session": adr_serverless.session,
        "dataset": adr_serverless.dataset,
    }
    objs = [
        HTML(name="test_create_objects_html", content="<h1>Heading 1</h1>", **kwargs),
        String(name="test_create_objects_string", content="This is a test string.", **kwargs),
    ]
    count = adr_serverless.create_objects(objs)
    assert count == 2, "No objects created"


@pytest.mark.ado_test
def test_create_objects_non_iter(adr_serverless):
    from ansys.dynamicreporting.core.serverless import HTML

    kwargs = {
        "session": adr_serverless.session,
        "dataset": adr_serverless.dataset,
    }
    obj = HTML(name="test_create_objects_html", content="<h1>Heading 1</h1>", **kwargs)
    with pytest.raises(ADRException, match="objects must be an iterable"):
        adr_serverless.create_objects(obj)


@pytest.mark.ado_test
def test_delete_items(adr_serverless):
    from ansys.dynamicreporting.core.serverless import File, Item

    # file
    adr_serverless.create_item(
        File,
        name="test_delete_items",
        content=str(Path(__file__).parent / "test_data" / "input.pptx"),
        tags="test_delete_items",
    )

    del_items = adr_serverless.query(query_type=Item, query="A|i_tags|cont|test_delete_items;")
    count = del_items.delete()
    assert count == 1, "No items deleted"


@pytest.mark.ado_test
def test_delete_sessions(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Session

    _ = Session.create(application="test_delete_sessions", tags="test_delete_sessions;")
    del_sesh = adr_serverless.query(query_type=Session, query="A|s_tags|cont|test_delete_sessions;")
    count = del_sesh.delete()
    assert count == 1, "No sessions deleted"


@pytest.mark.ado_test
def test_delete_datasets(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Dataset

    _ = Dataset.create(filename="test_delete_datasets", tags="test_delete_datasets;")
    del_dataset = adr_serverless.query(
        query_type=Dataset, query="A|d_tags|cont|test_delete_datasets"
    )
    count = del_dataset.delete()
    assert count == 1, "No datasets deleted"


@pytest.mark.ado_test
def test_delete_templates(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Template, TOCLayout

    _ = adr_serverless.create_template(TOCLayout, name="test_delete_templates")
    temps = adr_serverless.query(query_type=TOCLayout, query="A|t_name|eq|test_delete_templates;")
    count = temps.delete()
    assert count == 1, "No templates deleted"


@pytest.mark.ado_test
def test_render_report(monkeypatch, adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout
    from ansys.dynamicreporting.core.serverless import template as template_module

    # Create a dummy template so that Template.get() can find it.
    dummy_template = adr_serverless.create_template(
        BasicLayout, name="TestReport", parent=None, tags="dp=dp227"
    )
    dummy_template.set_params({"HTML": "<h1>Test Report</h1>"})
    dummy_template.save()

    # Patch django.template.loader.render_to_string to return a fixed string.
    def fake_render_to_string(template_name, context, request):
        # Check that the expected template file is requested.
        assert template_name == "reports/report_display_simple.html"
        # Instead of an exact match, check for a few essential keys.
        assert context.get("plotly") == 1 and context.get("page_width") == 10.5
        return "dummy rendered content"

    monkeypatch.setattr(template_module, "render_to_string", fake_render_to_string)

    # Call render_report; it looks up the template by name ("TestReport") and then calls render.
    result = adr_serverless.render_report(
        name="TestReport",
        context={"plotly": 1, "pwidth": "10.5", "dpi": "96."},
        item_filter="A|i_tags|cont|dp=dp227;",
    )

    # Assert that the result is the string returned by the monkey-patched render_to_string.
    assert result == "dummy rendered content"


@pytest.mark.ado_test
def test_render_no_kwarg(adr_serverless):
    with pytest.raises(ADRException, match="At least one keyword argument must be provided"):
        adr_serverless.render_report()


@pytest.mark.ado_test
def test_render_invalid_template(adr_serverless):
    with pytest.raises(ADRException, match="Report rendering failed"):
        adr_serverless.render_report(name="InvalidTemplateName")


@pytest.mark.ado_test
def test_render_report_as_pptx_success(adr_serverless, monkeypatch):
    from ansys.dynamicreporting.core.serverless import PPTXLayout
    from ansys.dynamicreporting.core.serverless import template as template_module

    # Create a valid PPTXLayout template
    _ = adr_serverless.create_template(PPTXLayout, name="TestPPTXReport", parent=None)

    # Mock the template's render_pptx method to avoid actual PPTX generation
    def fake_render_pptx(self, context, item_filter, request):
        return b"dummy pptx content"

    monkeypatch.setattr(template_module.PPTXLayout, "render_pptx", fake_render_pptx)

    # Call the method under test
    pptx_bytes = adr_serverless.render_report_as_pptx(
        name="TestPPTXReport", item_filter="A|i_tags|cont|dp=dp227;"
    )

    # Assert the result is the dummy content from the mock
    assert pptx_bytes == b"dummy pptx content"


@pytest.mark.ado_test
def test_render_report_as_pptx_no_kwarg(adr_serverless):
    with pytest.raises(ADRException, match="At least one keyword argument must be provided"):
        adr_serverless.render_report_as_pptx()


@pytest.mark.ado_test
def test_render_report_as_pptx_wrong_template_type(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout

    # Create a template that is NOT a PPTXLayout
    _ = adr_serverless.create_template(BasicLayout, name="NotAPPTXReport", parent=None)

    # Expect an ADRException because the template is not the correct type
    with pytest.raises(
        ADRException, match="The template must be of type 'PPTXLayout' to render as a PowerPoint"
    ):
        adr_serverless.render_report_as_pptx(name="NotAPPTXReport")


@pytest.mark.ado_test
def test_render_report_as_pptx_render_failure(adr_serverless, monkeypatch):
    from ansys.dynamicreporting.core.serverless import PPTXLayout
    from ansys.dynamicreporting.core.serverless import template as template_module

    # Create a valid PPTXLayout template
    _ = adr_serverless.create_template(PPTXLayout, name="FailingPPTXReport", parent=None)

    # Mock the underlying render_pptx method to simulate a failure
    def fake_render_pptx_fails(self, context, item_filter, request):
        raise Exception("Simulated rendering engine failure")

    monkeypatch.setattr(template_module.PPTXLayout, "render_pptx", fake_render_pptx_fails)

    # Expect an ADRException because the underlying render call failed
    with pytest.raises(ADRException, match="PPTX Report rendering failed"):
        adr_serverless.render_report_as_pptx(name="FailingPPTXReport")


@pytest.mark.ado_test
def test_full_pptx_report_generation_integration(adr_serverless):
    import datetime
    import io
    from pathlib import Path
    import random

    import numpy as np
    from pptx import Presentation

    from ansys.dynamicreporting.core.serverless import (
        HTML,
        File,
        Image,
        PPTXLayout,
        PPTXSlideLayout,
        Scene,
        String,
        Table,
        Tree,
    )

    source_tag = "pptx-test-serverless"  # A common tag to filter all items for this report

    # --- Source PPTX as a File item ---
    input_pptx_item_name = "input.pptx"
    adr_serverless.create_item(
        File,
        name=input_pptx_item_name,
        content=str(Path(__file__).parent / "test_data" / "input.pptx"),
        tags=source_tag,
    )

    # --- String items for various titles ---
    adr_serverless.create_item(
        String, name="title_text", content="My presentation", tags=source_tag
    )
    adr_serverless.create_item(
        String, name="toc_title", content="Table of contents", tags=source_tag
    )
    adr_serverless.create_item(String, name="toc_link_text", content="Go back", tags=source_tag)
    adr_serverless.create_item(String, name="html_title", content="My HTML item", tags=source_tag)
    adr_serverless.create_item(String, name="table_title", content="My table", tags=source_tag)
    adr_serverless.create_item(String, name="tree_title", content="My tree", tags=source_tag)
    adr_serverless.create_item(
        String, name="line_plot_title", content="My line plot", tags=source_tag
    )

    # --- HTML item ---
    html_item = adr_serverless.create_item(
        HTML,
        name="html",
        content=(
            "<h1>Heading 1</h1><h2>Heading 2</h2><h3>Heading 3</h3>"
            "<h4>Heading 4</h4><h5>Heading 5</h5>Two breaks below"
            "<br><br /><h6>Heading 6 (& one break below)</h6><br>The end"
        ),
        tags=f'{source_tag} pptx_slide_title="headers and breaks"',
    )

    # --- Image item ---
    image_item = adr_serverless.create_item(
        Image,
        name="logo",
        content=str(Path(__file__).parent / "test_data" / "nexus_logo.png"),
        tags=source_tag,
    )

    # --- Table items ---
    random.seed(12345)
    array1 = np.array(
        [
            [
                i,
                i * 200 - 1003,
                1.2 * i**3.4 + 3.5 * i + 123.0,
                (10 - i) ** 2.3,
                random.uniform(-2000, 6000),
            ]
            for i in range(10)
        ],
        dtype="f",
    )
    adr_serverless.create_item(
        Table,
        name="table1",
        content=array1,
        labels_col=["Linear", "Shift", "Polynomial", "Invert Poly", "Random"],
        title="Numeric table",
        tags=f'{source_tag} pptx_slide_title="<h2>Linear</h2>linear description<br /><br /><h4>Iterations: 10</h4>"',
    )

    random.seed(54321)
    array2 = np.array(
        [
            [
                i,
                i * 500 - 900,
                2.5 * i**3.4 + 5.5 * i + 65.0,
                (10 - i) ** 9.4,
                random.uniform(-5000, 3000),
            ]
            for i in range(10)
        ],
        dtype="f",
    )
    adr_serverless.create_item(
        Table,
        name="table2",
        content=array2,
        labels_col=["ID", "Location[X]", "Location[Y]", "Location[Z]", "turbViscosity[X]"],
        title="Numeric table2",
        tags=f'{source_tag} pptx_slide_title="Location-Viscosity"',
    )

    array3 = np.array(
        [['A {{"mylink"|nexus_link:"LINK"}}', "B \u4e14".encode(), "C"], [b"1", b"2", b"3"]],
        dtype="S50",
    )
    adr_serverless.create_item(
        Table,
        name="table3",
        content=array3,
        labels_row=["Row 1", "Row 2"],
        labels_col=["Column A", "Column B", "Column C"],
        title="Simple ASCII table",
        tags=source_tag,
    )

    # --- Tree item ---
    tree_content = [
        {
            "key": "root",
            "name": "Top Level",
            "value": None,
            "state": "collapsed",
            "tree_global_toggle": "1",
            "children": [
                {"key": "child", "name": "Boolean example", "value": True},
                {
                    "key": "child",
                    "name": "Simple string",
                    "value": 'Hello world!!! {{"mylink"|nexus_link:"LINK"}}',
                },
                {"key": "child", "name": "Integer example", "value": 10},
                {"key": "child", "name": "Float example", "value": 99.99},
                {"key": "child", "name": "multi-valued child", "value": ["val1", "val2"]},
                {"key": "child", "name": "Simple \u4e14 string", "value": "Hello \u4e14 world!!"},
                {"key": "child", "name": "Integer string 3", "value": "20200102"},
                {"key": "child", "name": "The current date", "value": datetime.datetime.now()},
                {"key": "child", "name": "A data item guid", "value": image_item.guid},
                {
                    "key": "child_parent",
                    "name": "A child parent",
                    "value": "Parents can have values",
                    "state": "expanded",
                    "children": [
                        {"key": "leaves", "name": "Leaf 0", "value": 0},
                        {"key": "leaves", "name": "Leaf 1", "value": 1},
                    ],
                },
            ],
        }
    ]
    adr_serverless.create_item(Tree, name="tree", content=tree_content, tags=source_tag)

    # --- Line Plot (as a Table item) ---
    line_plot_array = np.array(
        [
            [3.98, 4.41, 4.85, 5.29, 5.72, 6.16, 6.59, 7.03, 7.47, 7.90],
            [-5.08, -14.84, -24.19, -34.11, -45.64, -49.59, -52.44, -52.22, -50.30, -45.44],
        ],
        dtype="f",
    )
    adr_serverless.create_item(
        Table,
        name="line_plot",
        content=line_plot_array,
        labels_row=["X", "Lift"],
        title="Cumulative_Total_Lift",
        plot="line",
        tags=source_tag,
    )

    # --- ENS Session File Item---
    adr_serverless.create_item(
        File,
        name="session",
        content=str(Path(__file__).parent / "test_data" / "session.ens"),
        tags=f'{source_tag} pptx_slide_title="session-tag-title"',
    )

    # ==============================================================================
    # 2. Create the full template structure.
    # ==============================================================================
    report_name = "pptx-select"
    pptx_template = adr_serverless.create_template(PPTXLayout, name=report_name, parent=None)
    pptx_template.input_pptx = input_pptx_item_name
    pptx_template.output_pptx = "output-select.pptx"
    pptx_template.item_filter = f"A|i_src|cont|{source_tag};"
    pptx_template.use_all_slides = "0"
    pptx_template.save()

    # --- Define the slide children ---
    slides_to_create = [
        {"name": "start", "source_slide": "1", "filter": "A|i_name|eq|title_text;"},
        {"name": "toc", "source_slide": "2", "filter": "A|i_name|any|toc_title,toc_link_text;"},
        {"name": "html", "source_slide": "3", "filter": f"A|i_guid|eq|{html_item.guid};"},
        {
            "name": "table",
            "source_slide": "4",
            "properties": {"show_tag_title_only": "1"},
            "html": "<h1>table</h1>table description",
            "filter": "A|i_name|any|table1,table2,table3;",
        },
        {"name": "tree", "source_slide": "5", "filter": "A|i_name|cont|tree;"},
        {"name": "line", "source_slide": "6", "filter": "A|i_name|cont|line_plot;"},
        {
            "name": "session",
            "source_slide": "7",
            "filter": f"A|i_name|any|session;A|i_guid|eq|{image_item.guid};",
        },
    ]

    for slide_data in slides_to_create:
        slide = adr_serverless.create_template(
            PPTXSlideLayout, name=slide_data["name"], parent=pptx_template
        )
        slide.source_slide = slide_data["source_slide"]
        slide.item_filter = slide_data.get("filter", "")
        if "properties" in slide_data:
            slide.add_properties(slide_data["properties"])
        if "html" in slide_data:
            slide.set_html(slide_data["html"])
        slide.save()

    # ==============================================================================
    # 3.  Render the report.
    # ==============================================================================
    pptx_bytes = adr_serverless.render_report_as_pptx(name=report_name)

    # ==============================================================================
    # 4. Validate the output.
    # ==============================================================================
    assert isinstance(pptx_bytes, bytes)
    assert len(pptx_bytes) > 1000, "Generated PPTX file seems too small."

    try:
        pptx_file = io.BytesIO(pptx_bytes)
        prs = Presentation(pptx_file)

        # Expect 7 slides based on the template definition
        assert len(prs.slides) == 7

        # Spot check a few slides for expected content
        # Slide 1: Title
        title_slide_text = "".join(
            shape.text for shape in prs.slides[0].shapes if shape.has_text_frame
        )
        assert "My presentation" in title_slide_text

        # Slide 3: HTML content
        html_slide_text = "".join(
            shape.text for shape in prs.slides[2].shapes if shape.has_text_frame
        )
        assert "Heading 1" in html_slide_text and "The end" in html_slide_text

        # Slide 4: Table (check for title from HTML property)
        table_slide_text = "".join(
            shape.text for shape in prs.slides[3].shapes if shape.has_text_frame
        )
        assert "table description" in table_slide_text

    except Exception as e:
        pytest.fail(f"Failed to parse or validate the final PPTX file. Error: {e}")


@pytest.mark.ado_test
def test_copy_sessions(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Session

    tag = "dp=test_copy_sessions"

    Session.create(application="copy_sesh_1", tags=tag)
    Session.create(application="copy_sesh_2", tags=tag)

    count = adr_serverless.copy_objects(Session, "dest", query=f"A|s_tags|cont|{tag};")
    assert count == 2

    sessions = Session.filter(tags__icontains=tag, using="dest")
    assert len(sessions) == count


@pytest.mark.ado_test
def test_copy_datasets(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Dataset

    tag = "dp=test_copy_datasets"

    Dataset.create(filename="copy_dataset_1", tags=tag)
    Dataset.create(filename="copy_dataset_2", tags=tag)

    count = adr_serverless.copy_objects(Dataset, "dest", query=f"A|d_tags|cont|{tag};")
    assert count == 2

    datasets = Dataset.filter(tags__icontains=tag, using="dest")
    assert len(datasets) == count


@pytest.mark.ado_test
def test_copy_items(adr_serverless, tmp_path):
    from ansys.dynamicreporting.core.serverless import Image, Item, String

    tag = "dp=test_copy_items"

    adr_serverless.create_item(String, name="copy_item_1", content="This is a test item.", tags=tag)
    adr_serverless.create_item(
        Image,
        name="copy_item_2",
        content=str(Path(__file__).parent / "test_data" / "nexus_logo.png"),
        tags=tag,
    )

    count = adr_serverless.copy_objects(
        Item, "dest", query=f"A|i_tags|cont|{tag};", target_media_dir=tmp_path
    )
    assert count == 2

    items = Item.filter(tags__icontains=tag, using="dest")
    assert len(items) == count


@pytest.mark.ado_test
def test_copy_items_no_target_media_dir(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Image, Item

    tag = "dp=test_copy_items_no_target_media_dir"

    adr_serverless.create_item(
        Image,
        name="copy_item_4",
        content=str(Path(__file__).parent / "test_data" / "nexus_logo.png"),
        tags=tag,
    )

    count = adr_serverless.copy_objects(Image, "dest", query=f"A|i_tags|cont|{tag};")
    assert count == 1

    items = Item.filter(tags__icontains=tag, using="dest")
    assert len(items) == count


@pytest.mark.ado_test
def test_copy_items_test_run(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Image, Item

    tag = "dp=test_copy_items_test_run"
    adr_serverless.create_item(
        Image,
        name="test_copy_items_test_run",
        content=str(Path(__file__).parent / "test_data" / "nexus_logo.png"),
        tags=tag,
    )

    count = adr_serverless.copy_objects(Image, "dest", query=f"A|i_tags|cont|{tag};", test=True)
    assert count == 1

    items = Item.filter(tags__icontains=tag, using="dest")
    assert len(items) != count


@pytest.mark.ado_test
def test_copy_items_wrong_type(adr_serverless):
    tag = "dp=test_copy_items_wrong_type"
    with pytest.raises(TypeError, match="is not a type of"):
        adr_serverless.copy_objects(ADR, "dest", query=f"A|i_tags|cont|{tag};")


@pytest.mark.ado_test
def test_copy_items_invalid_database(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Item

    tag = "dp=test_copy_items_invalid_database"
    with pytest.raises(ADRException, match="must be configured first"):
        adr_serverless.copy_objects(Item, "invalid_db", query=f"A|i_tags|cont|{tag};")


@pytest.mark.ado_test
def test_copy_templates(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout, PanelLayout, Template

    tag = "dp=test_copy_templates"
    template_name = "test_copy_template_report"
    report = adr_serverless.create_template(BasicLayout, name=template_name, parent=None, tags=tag)
    adr_serverless.create_template(PanelLayout, name="Introduction", parent=report)

    count = adr_serverless.copy_objects(Template, "dest", query=f"A|t_name|eq|{template_name};")
    assert count == 1

    templates = Template.filter(name=template_name, using="dest")
    assert len(templates) == count


@pytest.mark.ado_test
def test_copy_templates_children(adr_serverless):
    from ansys.dynamicreporting.core.serverless import BasicLayout, PanelLayout, Template

    tag = "dp=test_copy_templates_children"
    template_name = "test_copy_templates_children_report"
    report = adr_serverless.create_template(BasicLayout, name=template_name, parent=None, tags=tag)
    adr_serverless.create_template(PanelLayout, name="Introduction", parent=report, tags=tag)

    with pytest.raises(ADRException, match="Only top-level templates can be copied"):
        adr_serverless.copy_objects(Template, "dest", query=f"A|t_tags|cont|{tag};")


@pytest.mark.ado_test
def test_load_templates_from_file(adr_serverless):
    from ansys.dynamicreporting.core.serverless import Template

    # Load templates from the sample JSON file
    sample_file = Path(__file__).parent.parent / "test_data" / "sample.json"
    adr_serverless.load_templates_from_file(sample_file)

    # Verify the root template
    root_template = (adr_serverless.query(query_type=Template, query="A|t_name|eq|A;"))[0]
    # root_template.store_json("haha.json")
    assert root_template is not None
    assert root_template.name == "A"
    assert root_template.report_type == "Layout:basic"
    assert root_template.get_params()["HTML"] == "<h1>Serverless Simulation Report</h1>"

    # Verify child templates
    child_templates = root_template.children
    assert len(child_templates) == 2

    child_b = next((child for child in child_templates if child.name == "B"), None)
    assert child_b is not None
    assert child_b.report_type == "Layout:panel"

    child_c = next((child for child in child_templates if child.name == "C"), None)
    assert child_c is not None
    assert child_c.report_type == "Layout:basic"
    assert child_c.get_params()["HTML"] == "<h2>Basic C</h2>"

    # Verify grandchild template
    grandchild_d = next((child for child in child_c.children if child.name == "D"), None)
    assert grandchild_d is not None
    assert grandchild_d.report_type == "Layout:basic"
    assert grandchild_d.get_params()["HTML"] == "<h2>Basic D</h2>"


@pytest.mark.ado_test
def test_load_templates_from_file_no_such_file(adr_serverless):
    with pytest.raises(FileNotFoundError, match="The file 'nonexistent.json' does not exist."):
        adr_serverless.load_templates_from_file("nonexistent.json")
